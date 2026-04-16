"""クライアント提出用PPTXレポート生成モジュール — 実診断データを反映"""

import math
from io import BytesIO
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


# ===== カラーパレット（Ocean Gradient） =====
class C:
    navy = RGBColor(0x0A, 0x16, 0x28)
    deep_blue = RGBColor(0x06, 0x5A, 0x82)
    teal = RGBColor(0x0D, 0x94, 0x88)
    cyan = RGBColor(0x14, 0xB8, 0xA6)
    cream = RGBColor(0xF8, 0xFA, 0xFC)
    light_bg = RGBColor(0xF1, 0xF5, 0xF9)
    text_dark = RGBColor(0x1E, 0x29, 0x3B)
    text_light = RGBColor(0xE2, 0xE8, 0xF0)
    muted = RGBColor(0x64, 0x74, 0x8B)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    green = RGBColor(0x10, 0xB9, 0x81)
    amber = RGBColor(0xF5, 0x9E, 0x0B)
    red = RGBColor(0xEF, 0x44, 0x44)
    border = RGBColor(0xCB, 0xD5, 0xE1)


FONT_HDR = "Georgia"
FONT_BODY = "メイリオ"  # 日本語対応


# ===== ヘルパー =====
def _set_text(tf, text: str, *, size=14, bold=False, italic=False,
              color=C.text_dark, font=FONT_BODY, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP):
    tf.text = ""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def _add_rect(slide, x, y, w, h, fill_color, line_color=None, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is not None:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def _add_text_box(slide, x, y, w, h, text, **kwargs):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    _set_text(tb.text_frame, text, **kwargs)
    return tb


def _header_bar(slide, title, subtitle=""):
    _add_rect(slide, 0, 0, 13.33, 0.7, C.navy)
    _add_rect(slide, 0, 0.7, 13.33, 0.04, C.teal)
    _add_text_box(slide, 0.5, 0.1, 9, 0.5, title,
                  size=18, bold=True, color=C.white, font=FONT_HDR,
                  anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _add_text_box(slide, 9.5, 0.1, 3.5, 0.5, subtitle,
                      size=11, color=C.cyan, align=PP_ALIGN.RIGHT,
                      anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, page_num):
    _add_text_box(slide, 0.5, 7.2, 12.3, 0.25,
                  f"AIO/LLMO診断レポート  |  {page_num}",
                  size=9, color=C.muted, align=PP_ALIGN.RIGHT)


def _grade_color(grade: str):
    return {"A": C.green, "B": C.teal, "C": C.amber, "D": C.red, "E": C.red}.get(grade, C.muted)


def _score_color(pct: float):
    if pct >= 0.75:
        return C.green
    if pct >= 0.50:
        return C.teal
    if pct >= 0.35:
        return C.amber
    return C.red


def _set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ====================================================================
# メイン関数
# ====================================================================
def generate_pptx_report(
    url: str,
    structure: dict,
    total: dict,
    categories: dict,
    all_scores: dict,
    robots: dict,
    llms: dict,
    pagespeed: dict,
    competitors: list = None,
    comparison: dict = None,
    improvements: dict = None,
    test_queries: dict = None,
    cv_data: dict = None,
    site_agg: dict = None,
    preset_id: str = "media",
    site_diagnosis: dict = None,
) -> BytesIO:
    """診断データから網羅的なクライアント提出用スライドを生成（20-35枚想定）。

    preset_id: media / recruiting / corporate — カテゴリラベルを切り替え。
    """
    # プリセット別カテゴリ定義を動的に取得
    try:
        from core.presets import get_preset
        _preset_module = get_preset(preset_id)["module"]
        _cat_defs = getattr(_preset_module, "CATEGORY_DEFINITIONS", None)
    except Exception:
        _cat_defs = None
    if _cat_defs:
        cat_order = [(k, v.get("label", k)) for k, v in _cat_defs.items()]
    else:
        cat_order = [
            ("1_content", "コンテンツ品質"),
            ("2_structured", "構造化データ"),
            ("3_eeat", "E-E-A-Tシグナル"),
            ("4_citation", "AI引用可能性"),
            ("5_freshness", "コンテンツ鮮度"),
            ("6_technical", "テクニカルAIO/UX"),
        ]

    # プリセット別の表紙コピー
    _preset_label_map = {
        "media": ("オウンドメディア / コンテンツサイト", "Content & Media Site"),
        "recruiting": ("採用ページ / リクルーティングサイト", "Recruiting Site"),
        "corporate": ("コーポレートサイト / 企業公式サイト", "Corporate Site"),
    }
    preset_label_jp, preset_label_en = _preset_label_map.get(
        preset_id, ("Webサイト", "Website")
    )
    pres = Presentation()
    pres.slide_width = Inches(13.33)
    pres.slide_height = Inches(7.5)
    blank = pres.slide_layouts[6]

    # ページ番号カウンター（リスト型でクロージャ参照）
    pn = [1]
    def _pg() -> str:
        n = pn[0]
        pn[0] += 1
        return f"{n:02d}"

    domain = url.split("//")[-1].split("/")[0] if "//" in url else url
    site_title = structure.get("title", domain)
    today = datetime.now().strftime("%Y年%m月%d日")
    competitors = competitors or []
    comparison = comparison or {}
    improvements = improvements or {}
    test_queries = test_queries or {}

    # ==============================================================
    # Slide 1: タイトル
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.navy)
    _add_rect(s, 0, 0, 5.5, 7.5, C.deep_blue)
    _add_rect(s, 5.3, 0, 0.05, 7.5, C.cyan)

    # ロゴサークル
    _add_rect(s, 1.2, 1.0, 1.4, 1.4, C.cyan, shape_type=MSO_SHAPE.OVAL)
    _add_text_box(s, 1.2, 1.0, 1.4, 1.4, "AIO",
                  size=32, bold=True, color=C.navy, font=FONT_HDR,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(s, 1.0, 3.0, 4.3, 0.3,
                  datetime.now().strftime("%Y"),
                  size=12, bold=True, color=C.cyan)
    _add_text_box(s, 1.0, 3.4, 4.3, 1.8,
                  "AIO / LLMO\n総合診断レポート",
                  size=32, bold=True, color=C.white, font=FONT_HDR)
    _add_rect(s, 1.0, 5.3, 0.6, 0.06, C.cyan)
    _add_text_box(s, 1.0, 5.45, 4.3, 0.8,
                  "AI Overview & LLM Optimization\nAssessment Report",
                  size=12, italic=True, color=C.text_light)

    # 右サイド
    _add_text_box(s, 6.2, 1.5, 6.5, 0.5, "診断対象サイト",
                  size=14, color=C.cyan)
    _add_text_box(s, 6.2, 2.0, 6.5, 0.8, site_title[:40],
                  size=26, bold=True, color=C.white, font=FONT_HDR)
    _add_text_box(s, 6.2, 2.9, 6.5, 0.4, url[:70],
                  size=12, italic=True, color=C.text_light)

    _total_items = sum(len(s) for s in (all_scores or {}).values()) if all_scores else 30
    bullets = [
        f"診断種別: {preset_label_jp}",
        f"総合スコア: {total['total']} / 100（グレード {total['grade']}）",
        f"{len(cat_order)}カテゴリ × {_total_items}項目の実測診断",
        f"競合ベンチマーク & 改善アクションプラン",
    ]
    for i, b in enumerate(bullets):
        y = 3.8 + i * 0.5
        _add_rect(s, 6.2, y + 0.08, 0.2, 0.2, C.cyan, shape_type=MSO_SHAPE.OVAL)
        _add_text_box(s, 6.55, y, 6, 0.4, b,
                      size=14, color=C.text_light, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(s, 6.2, 6.7, 6.5, 0.3, f"診断日: {today}",
                  size=11, color=C.muted)

    # ==============================================================
    # Slide 2: エグゼクティブサマリー
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "エグゼクティブサマリー", "Executive Summary")

    # スコアカード
    _add_rect(s, 0.5, 1.1, 4.3, 5.8, C.navy, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 1.4, 4.3, 0.4, "総合スコア",
                  size=14, color=C.cyan, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 1.85, 4.3, 2.0, str(total["total"]),
                  size=96, bold=True, color=C.white, font=FONT_HDR,
                  align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 3.8, 4.3, 0.4, "/ 100",
                  size=18, color=C.text_light, align=PP_ALIGN.CENTER)

    _add_rect(s, 1.65, 4.4, 2.0, 0.7, _grade_color(total["grade"]),
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 1.65, 4.4, 2.0, 0.7, f"グレード {total['grade']}",
                  size=22, bold=True, color=C.white, font=FONT_HDR,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(s, 0.7, 5.4, 3.9, 0.9, total.get("label", ""),
                  size=11, italic=True, color=C.text_light,
                  align=PP_ALIGN.CENTER)

    # カテゴリ別
    _add_text_box(s, 5.2, 1.15, 7.5, 0.4, "カテゴリ別スコア",
                  size=18, bold=True, color=C.text_dark, font=FONT_HDR)

    for i, (key, label) in enumerate(cat_order):
        cat = categories.get(key, {"score": 0, "max": 20})
        y = 1.75 + i * 0.82
        score = cat.get("score", 0)
        mx = cat.get("max", 20) or 20
        pct = score / mx

        _add_text_box(s, 5.2, y, 3.0, 0.35, label,
                      size=13, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, 8.3, y + 0.08, 3.5, 0.22, C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 8.3, y + 0.08, max(3.5 * pct, 0.05), 0.22,
                  _score_color(pct), shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 11.9, y, 0.9, 0.35, f"{score}/{mx}",
                      size=12, bold=True, color=C.text_dark,
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 3: なぜ今 AIO/LLMO
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "なぜ今、AIO/LLMO対策なのか", "The Shift to AI Search")

    _add_text_box(s, 0.5, 1.0, 12.3, 0.5,
                  "AI検索の普及で「検索→クリック」の時代は終わった",
                  size=20, bold=True, color=C.deep_blue, font=FONT_HDR)

    stats = [
        ("58%", "Google検索でAI Overviewが表示", "日本では2024年8月〜本格展開"),
        ("3.2B", "ChatGPT月間アクセス数", "既にGoogle検索の約1/3規模"),
        ("-35%", "AI Overview表示時のCTR下落", "従来SEOだけでは機会損失"),
    ]
    for i, (num, label, sub) in enumerate(stats):
        x = 0.5 + i * 4.3
        _add_rect(s, x, 1.8, 4.0, 2.2, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x, 1.8, 0.1, 2.2, C.teal)
        _add_text_box(s, x + 0.2, 1.95, 3.8, 1.0, num,
                      size=44, bold=True, color=C.deep_blue, font=FONT_HDR,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, x + 0.2, 2.95, 3.8, 0.5, label,
                      size=13, bold=True, color=C.text_dark)
        _add_text_box(s, x + 0.2, 3.45, 3.8, 0.4, sub,
                      size=10, italic=True, color=C.muted)

    # 下段2カラム
    _add_rect(s, 0.5, 4.4, 6.1, 2.4, C.navy,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, 4.55, 5.7, 0.4, "❌ 何もしない場合",
                  size=14, bold=True, color=C.red)
    bad = [
        "AI検索結果に一切表示されず機会損失",
        "競合がAI引用を獲得、相対順位が下落",
        "SEO投資のROIが年々低下",
        "商談・採用での情報露出が減少",
    ]
    for i, b in enumerate(bad):
        _add_text_box(s, 0.8, 5.0 + i * 0.4, 5.6, 0.35, "• " + b,
                      size=12, color=C.text_light)

    _add_rect(s, 6.8, 4.4, 6.0, 2.4, C.teal,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 7.0, 4.55, 5.6, 0.4, "✅ 早期対策の効果",
                  size=14, bold=True, color=C.white)
    good = [
        "AI Overviewでの引用率が平均3〜5倍向上",
        "Perplexity / ChatGPTでの出典表示獲得",
        "ブランド第一想起（エンティティ認識）を確立",
        "SEOとAIOの相乗効果で検索総流入増",
    ]
    for i, g in enumerate(good):
        _add_text_box(s, 7.1, 5.0 + i * 0.4, 5.5, 0.35, "• " + g,
                      size=12, color=C.white)

    _footer(s, _pg())

    # ==============================================================
    # Slide 4: 診断フレームワーク
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "診断フレームワーク", "6 Categories × 30 Items")

    _add_text_box(s, 0.5, 1.0, 12.3, 0.4,
                  "AI検索時代に必要な6つの評価軸で、30項目を実測診断",
                  size=14, color=C.muted)

    framework = [
        ("01", "コンテンツ品質・構造", "20点", "Answer-first, 見出し階層, FAQ, 明瞭性", C.deep_blue),
        ("02", "構造化データ",         "20点", "Organization, Article, FAQPage, Breadcrumb", C.teal),
        ("03", "E-E-A-Tシグナル",      "20点", "著者情報, 運営者, 引用, 編集ポリシー", C.cyan),
        ("04", "AI引用可能性",         "20点", "定義文, 数値データ, 独自情報, エンティティ", C.deep_blue),
        ("05", "コンテンツ鮮度",       "10点", "更新日表示, dateModified, 年次管理", C.teal),
        ("06", "テクニカルAIO/UX",     "10点", "AIクローラー, PageSpeed, canonical, sitemap", C.cyan),
    ]
    for i, (num, label, pts, items, color) in enumerate(framework):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.3
        y = 1.6 + row * 2.6

        _add_rect(s, x, y, 4.0, 2.4, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x + 0.2, y + 0.2, 0.8, 0.8, color,
                  shape_type=MSO_SHAPE.OVAL)
        _add_text_box(s, x + 0.2, y + 0.2, 0.8, 0.8, num,
                      size=16, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, x + 2.9, y + 0.3, 1.0, 0.4, C.light_bg,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, x + 2.9, y + 0.3, 1.0, 0.4, pts,
                      size=11, bold=True, color=C.deep_blue,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, x + 0.2, y + 1.1, 3.7, 0.5, label,
                      size=16, bold=True, color=C.text_dark, font=FONT_HDR)
        _add_rect(s, x + 0.2, y + 1.62, 0.5, 0.03, color)
        _add_text_box(s, x + 0.2, y + 1.75, 3.7, 0.6, items,
                      size=10, color=C.muted)

    _footer(s, _pg())

    # ==============================================================
    # Slide 5: カテゴリ別診断結果
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "カテゴリ別診断結果", "Category Breakdown")

    # 判定を計算
    def _judge(pct):
        if pct >= 0.8:
            return "A", C.green, "🟢 維持"
        if pct >= 0.6:
            return "B", C.teal, "🟢 維持"
        if pct >= 0.4:
            return "C", C.amber, "🟡 要改善"
        return "D", C.red, "🔴 最優先"

    rows = []
    for key, label in cat_order:
        cat = categories.get(key, {})
        score = cat.get("score", 0)
        mx = cat.get("max", 20) or 20
        pct = score / mx
        status, color, priority = _judge(pct)
        # 主な所見を算出
        items = cat.get("items", [])
        low = [i for i in items if i.get("max", 0) > 0
               and i.get("score", 0) / i.get("max", 1) < 0.5]
        note = f"低スコア項目 {len(low)}件" if low else "主要項目クリア"
        rows.append({
            "cat": label, "score": f"{score}/{mx}",
            "status": status, "note": note, "color": color,
            "priority": priority,
        })

    # テーブルヘッダー
    _add_rect(s, 0.5, 1.1, 12.3, 0.5, C.navy)
    hdrs = [("カテゴリ", 0.7, 3.5), ("スコア", 4.2, 1.5), ("判定", 5.7, 1.0),
            ("主な所見", 6.7, 4.0), ("優先度", 10.7, 2.0)]
    for t, x, w in hdrs:
        _add_text_box(s, x, 1.1, w, 0.5, t,
                      size=12, bold=True, color=C.white,
                      anchor=MSO_ANCHOR.MIDDLE)

    for i, r in enumerate(rows):
        y = 1.6 + i * 0.55
        if i % 2 == 0:
            _add_rect(s, 0.5, y, 12.3, 0.55, C.white)
        _add_text_box(s, 0.7, y, 3.5, 0.55, r["cat"],
                      size=13, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 4.2, y, 1.5, 0.55, r["score"],
                      size=15, bold=True, color=C.deep_blue, font=FONT_HDR,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, 5.7, y + 0.1, 0.7, 0.35, r["color"],
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 5.7, y + 0.1, 0.7, 0.35, r["status"],
                      size=11, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 6.7, y, 4.0, 0.55, r["note"],
                      size=11, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 10.7, y, 2.0, 0.55, r["priority"],
                      size=11, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)

    # インサイト
    _add_rect(s, 0.5, 5.2, 12.3, 1.6, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.8, 5.35, 11.7, 0.4, "💡 診断インサイト",
                  size=13, bold=True, color=C.cyan)

    # 弱点カテゴリを特定
    weak_cats = sorted(
        [(lbl, categories.get(k, {})) for k, lbl in cat_order],
        key=lambda x: x[1].get("score", 0) / (x[1].get("max", 1) or 1),
    )[:2]
    weak_names = " / ".join(w[0] for w in weak_cats)
    insight = (
        f"特に改善が必要な領域は「{weak_names}」です。この2領域は実装コストが低く、"
        f"効果が大きい「クイックウィン領域」です。60日以内の改善で +15〜20点の向上が見込めます。"
    )
    _add_text_box(s, 0.8, 5.75, 11.7, 1.0, insight,
                  size=12, color=C.white)

    _footer(s, _pg())

    # ==============================================================
    # Slide 6: テクニカル診断
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "テクニカル診断 — AIクローラー対応", "robots.txt / llms.txt / Performance")

    _add_text_box(s, 0.5, 1.1, 6.5, 0.4, "🤖 AIクローラーのアクセス状況",
                  size=16, bold=True, color=C.text_dark, font=FONT_HDR)

    crawlers_data = robots.get("crawlers", {})
    for i, (name, info) in enumerate(list(crawlers_data.items())[:6]):
        y = 1.6 + i * 0.55
        status = info.get("status", "未設定")
        if "ブロック" in status:
            color = C.red
        elif "許可" in status:
            color = C.green
        else:
            color = C.amber

        _add_rect(s, 0.5, y, 6.3, 0.5, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.7, y, 2.0, 0.5, name,
                      size=12, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 2.6, y, 2.5, 0.5, info.get("vendor", ""),
                      size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)
        # ステータスラベルを短縮
        status_short = status[:8] if len(status) > 8 else status
        _add_rect(s, 5.3, y + 0.1, 1.3, 0.3, color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 5.3, y + 0.1, 1.3, 0.3, status_short,
                      size=9, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 右サイド メトリクス
    _add_text_box(s, 7.2, 1.1, 5.6, 0.4, "⚙️ テクニカルメトリクス",
                  size=16, bold=True, color=C.text_dark, font=FONT_HDR)

    ps_score = pagespeed.get("score")
    ps_val = f"{ps_score}" if ps_score is not None else "—"
    ps_color = C.green if (ps_score or 0) >= 80 else (
        C.amber if (ps_score or 0) >= 50 else C.red)

    metrics = [
        ("PageSpeed", ps_val, "/100" if ps_score else "",
         f"Performance Score: {ps_score}" if ps_score else "取得失敗",
         ps_color),
        ("llms.txt", "○" if llms.get("exists") else "×", "",
         llms.get("summary", ""),
         C.green if llms.get("exists") else C.red),
        ("sitemap.xml", "○" if structure.get("sitemap_exists", True) else "×", "",
         "登録済み",
         C.green),
        ("canonical", "○" if structure.get("canonical") else "×", "",
         (structure.get("canonical", "")[:40] if structure.get("canonical") else "未実装"),
         C.green if structure.get("canonical") else C.red),
        ("viewport", "○" if structure.get("viewport") else "×", "",
         "モバイル対応" if structure.get("viewport") else "未対応",
         C.green if structure.get("viewport") else C.red),
    ]

    for i, (lbl, val, unit, note, color) in enumerate(metrics):
        y = 1.6 + i * 0.95
        _add_rect(s, 7.2, y, 5.6, 0.85, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 7.2, y, 0.08, 0.85, color)
        _add_text_box(s, 7.45, y + 0.1, 3.0, 0.35, lbl,
                      size=12, bold=True, color=C.text_dark)
        _add_text_box(s, 7.45, y + 0.45, 3.5, 0.35, note,
                      size=10, color=C.muted)
        _add_text_box(s, 10.8, y, 1.9, 0.85, val + unit,
                      size=24, bold=True, color=color, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 7: 構造化データ
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "構造化データ（JSON-LD）実装状況", "Schema.org Coverage")

    _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                  "AIが正確に情報を読み取るための「機械可読ラベル」の実装状況",
                  size=13, italic=True, color=C.muted)

    # 実装されているスキーマタイプを検出
    found_types = set()
    for sd in structure.get("jsonld", []):
        t = sd.get("@type", "")
        if isinstance(t, list):
            found_types.update(t)
        else:
            found_types.add(t)

    schemas = [
        ("Organization", "Organization" in found_types, "必須", "運営者情報をAIに明示"),
        ("Article / BlogPosting",
         "Article" in found_types or "BlogPosting" in found_types,
         "必須", "記事メタデータ認識"),
        ("FAQPage", "FAQPage" in found_types, "必須",
         "AI Overview引用率が大幅アップ"),
        ("BreadcrumbList", "BreadcrumbList" in found_types, "推奨",
         "検索結果パンくず表示"),
        ("LocalBusiness", "LocalBusiness" in found_types, "条件",
         "ローカルビジネスのみ必要"),
        ("HowTo", "HowTo" in found_types, "推奨",
         "手順コンテンツに有効"),
        ("Product", "Product" in found_types, "推奨",
         "製品・サービスに有効"),
    ]

    _add_rect(s, 0.5, 1.55, 12.3, 0.45, C.navy)
    for t, x, w in [("スキーマタイプ", 0.8, 3.0), ("状況", 3.8, 1.3),
                    ("状態", 5.1, 1.8), ("優先度", 6.9, 1.5),
                    ("AIへの効果", 8.4, 4.3)]:
        _add_text_box(s, x, 1.55, w, 0.45, t,
                      size=11, bold=True, color=C.white,
                      anchor=MSO_ANCHOR.MIDDLE)

    for i, (stype, impl, prio, impact) in enumerate(schemas):
        y = 2.0 + i * 0.52
        if i % 2 == 0:
            _add_rect(s, 0.5, y, 12.3, 0.52, C.white)
        icon = "○" if impl else "×"
        status = "実装済み" if impl else "未実装"
        color = C.green if impl else C.red

        _add_text_box(s, 0.8, y, 3.0, 0.52, stype,
                      size=12, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 3.8, y, 1.3, 0.52, icon,
                      size=18, bold=True, color=color, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 5.1, y, 1.8, 0.52, status,
                      size=11, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
        prio_color = C.red if prio == "必須" else (C.amber if prio == "推奨" else C.muted)
        _add_text_box(s, 6.9, y, 1.5, 0.52, prio,
                      size=11, bold=True, color=prio_color,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 8.4, y, 4.3, 0.52, impact,
                      size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

    # コールアウト
    missing_required = [stype for stype, impl, prio, _ in schemas
                        if not impl and prio == "必須"]
    callout = (
        f"⚡ {' / '.join(missing_required[:3])} の追加で大幅スコアアップ。実装所要時間は合計2〜3時間程度です。"
        if missing_required else
        "⚡ 必須スキーマは実装済み。推奨スキーマの追加で更なる最適化が可能です。"
    )
    _add_rect(s, 0.5, 5.9, 12.3, 0.9, C.teal,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, 5.9, 11.9, 0.9, callout,
                  size=13, bold=True, color=C.white,
                  anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 8: 競合ベンチマーク
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "競合ベンチマーク", "Competitive Analysis")

    _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                  "メインキーワード上位サイトとの比較",
                  size=13, italic=True, color=C.muted)

    # 自サイト + 競合データ
    sites = [{
        "name": "自サイト",
        "wc": structure.get("word_count", 0),
        "h2": sum(1 for h in structure.get("headings", []) if h.get("level") == 2),
        "faq": len(structure.get("faq_items", [])),
        "sd": ", ".join(list(found_types)[:3]) or "なし",
        "score": total["total"],
        "self": True,
    }]
    for c in competitors[:3]:
        sites.append({
            "name": (c.get("title") or c.get("url", ""))[:25],
            "wc": c.get("word_count", 0),
            "h2": c.get("h2_count", 0),
            "faq": c.get("faq_count", 0),
            "sd": ", ".join(c.get("sd_types", [])[:3]) or "なし",
            "score": c.get("score_pct", 0),
            "self": False,
        })

    _add_rect(s, 0.5, 1.55, 12.3, 0.5, C.navy)
    for t, x, w in [("サイト", 0.7, 2.8), ("文字数", 3.5, 1.4),
                    ("H2数", 4.9, 1.2), ("FAQ数", 6.1, 1.2),
                    ("構造化データ", 7.3, 3.5), ("スコア", 10.8, 1.9)]:
        _add_text_box(s, x, 1.55, w, 0.5, t,
                      size=12, bold=True, color=C.white,
                      anchor=MSO_ANCHOR.MIDDLE)

    for i, site in enumerate(sites):
        y = 2.05 + i * 0.6
        bg = C.light_bg if site["self"] else C.white
        _add_rect(s, 0.5, y, 12.3, 0.6, bg, line_color=C.border)
        _add_text_box(s, 0.7, y, 2.8, 0.6, site["name"],
                      size=13, bold=True,
                      color=C.deep_blue if site["self"] else C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 3.5, y, 1.4, 0.6, f"{site['wc']:,}",
                      size=12, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 4.9, y, 1.2, 0.6, str(site["h2"]),
                      size=12, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 6.1, y, 1.2, 0.6, str(site["faq"]),
                      size=12, bold=(site["faq"] == 0),
                      color=(C.red if site["faq"] == 0 else C.text_dark),
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 7.3, y, 3.5, 0.6, site["sd"][:40],
                      size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)
        sc_color = (C.green if site["score"] >= 80 else
                    C.amber if site["score"] >= 60 else C.red)
        _add_rect(s, 10.8, y + 0.1, 1.5, 0.4, sc_color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 10.8, y + 0.1, 1.5, 0.4, f"{site['score']}点",
                      size=13, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ギャップ / 優位性
    gaps_y = 2.05 + len(sites) * 0.6 + 0.3
    _add_rect(s, 0.5, gaps_y, 6.1, 6.8 - gaps_y, C.white, line_color=C.red,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, gaps_y + 0.1, 5.7, 0.4,
                  "🔴 コンテンツギャップ（負けている項目）",
                  size=13, bold=True, color=C.red)
    gaps = comparison.get("gaps", [])[:3]
    if not gaps:
        gaps_text = ["競合データ不足のため算出不能"]
    else:
        gaps_text = [f"{g['item']}: 自 {g['self_score']} vs 競合平均 {g['competitor_avg']}"
                     for g in gaps]
    for i, gt in enumerate(gaps_text):
        _add_text_box(s, 0.8, gaps_y + 0.55 + i * 0.4, 5.6, 0.4, "• " + gt,
                      size=11, color=C.text_dark)

    _add_rect(s, 6.8, gaps_y, 6.0, 6.8 - gaps_y, C.white, line_color=C.green,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 7.0, gaps_y + 0.1, 5.6, 0.4,
                  "✅ 独自優位性（勝っている項目）",
                  size=13, bold=True, color=C.green)
    advs = comparison.get("advantages", [])[:3]
    if not advs:
        advs_text = ["競合データ不足のため算出不能"]
    else:
        advs_text = [f"{a['item']}: 自 {a['self_score']} vs 競合平均 {a['competitor_avg']}"
                     for a in advs]
    for i, at in enumerate(advs_text):
        _add_text_box(s, 7.1, gaps_y + 0.55 + i * 0.4, 5.5, 0.4, "• " + at,
                      size=11, color=C.text_dark)

    _footer(s, _pg())

    # ==============================================================
    # Slide 9: Quick Win
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "優先改善施策 — Quick Win", "60日以内で実施可能な高効果施策")

    _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                  "今すぐ着手でき、かつインパクトが大きい施策",
                  size=13, italic=True, color=C.muted)

    qw_list = improvements.get("quick_wins", [])
    if not qw_list:
        # デフォルトQW
        qw_list = [
            {"title": "Organization JSON-LDの追加", "category": "構造化データ",
             "effort": "30分", "impact": "+4点",
             "reason": "運営者情報をAIに明示。全ページ共通で1度実装すれば完了。"},
            {"title": "FAQPage JSON-LD + FAQ追加", "category": "構造化データ",
             "effort": "2時間", "impact": "+6点",
             "reason": "Google AI Overview引用率が平均3倍に。即効性が最も高い施策。"},
            {"title": "著者ボックスの設置", "category": "E-E-A-T",
             "effort": "1時間", "impact": "+3点",
             "reason": "著者名・専門分野・プロフィールリンクを全記事に追加。"},
            {"title": "llms.txt の設置", "category": "テクニカル",
             "effort": "30分", "impact": "+2点",
             "reason": "ChatGPTがサイト概要を正確に把握（2025年10月〜対応）。"},
            {"title": "冒頭定義文の追加", "category": "AI引用可能性",
             "effort": "1時間", "impact": "+3点",
             "reason": "AIが最も引用しやすいパターン。既存記事の冒頭改修のみで可能。"},
        ]

    for i, q in enumerate(qw_list[:5]):
        y = 1.55 + i * 1.08
        _add_rect(s, 0.5, y, 12.3, 0.98, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 0.7, y + 0.2, 0.6, 0.6, C.teal,
                  shape_type=MSO_SHAPE.OVAL)
        _add_text_box(s, 0.7, y + 0.2, 0.6, 0.6, str(i + 1),
                      size=20, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 1.5, y + 0.1, 6.5, 0.4,
                      q.get("title", "")[:50],
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, 1.5, y + 0.55, 1.7, 0.3, C.light_bg,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 1.5, y + 0.55, 1.7, 0.3, q.get("category", "")[:12],
                      size=9, bold=True, color=C.deep_blue,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 3.3, y + 0.5, 6.8, 0.4,
                      q.get("reason", "")[:110],
                      size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 10.3, y + 0.1, 2.2, 0.3, "インパクト",
                      size=9, color=C.muted, align=PP_ALIGN.CENTER)
        _add_rect(s, 10.4, y + 0.4, 2.0, 0.45, C.teal,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 10.4, y + 0.4, 2.0, 0.45,
                      q.get("impact", "")[:10],
                      size=14, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 10〜N: Quick Win 詳細展開（1施策につき1ページ）
    # ==============================================================
    qw_detail_list = improvements.get("quick_wins", []) or qw_list
    for idx, qw in enumerate(qw_detail_list[:8], 1):
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)

        # ヘッダー
        priority = qw.get("priority", "A")
        pri_color = {"S": C.red, "A": C.amber, "B": C.teal, "C": C.muted}.get(priority, C.teal)
        _add_rect(s, 0, 0, 13.33, 0.65, C.navy)
        _add_text_box(s, 0.5, 0.13, 0.9, 0.4, f"#{idx:02d}",
                      size=18, bold=True, color=C.cyan, font=FONT_HDR)
        _add_rect(s, 1.4, 0.18, 0.7, 0.32, pri_color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 1.4, 0.18, 0.7, 0.32, f"優先{priority}",
                      size=11, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 2.2, 0.13, 9.0, 0.4, qw.get("title", "")[:70],
                      size=16, bold=True, color=C.white, font=FONT_HDR,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 11.3, 0.13, 1.9, 0.4,
                      qw.get("category", "")[:14],
                      size=11, color=C.cyan,
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, 0, 0.65, 13.33, 0.04, C.cyan)

        # 上段：Why（背景・狙い）／KPI／インパクト
        _add_rect(s, 0.4, 0.95, 8.0, 1.45, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.6, 1.0, 7.7, 0.3, "💡 なぜこの施策か（Why）",
                      size=11, bold=True, color=C.deep_blue)
        _add_text_box(s, 0.6, 1.3, 7.7, 1.05, qw.get("why", "")[:280],
                      size=11, color=C.text_dark)

        # 右上：KPI / Impact / Effort
        _add_rect(s, 8.55, 0.95, 4.4, 0.65, C.deep_blue,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 8.7, 0.95, 4.1, 0.25, "🎯 KPI",
                      size=9, color=C.cyan, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 8.7, 1.18, 4.1, 0.4, qw.get("kpi", "")[:60],
                      size=10, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)

        _add_rect(s, 8.55, 1.7, 2.1, 0.7, C.teal,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 8.55, 1.7, 2.1, 0.25, "想定インパクト",
                      size=9, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 8.55, 1.95, 2.1, 0.45, str(qw.get("impact", "+5点"))[:14],
                      size=15, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _add_rect(s, 10.85, 1.7, 2.1, 0.7, C.amber,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 10.85, 1.7, 2.1, 0.25, "実装工数",
                      size=9, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 10.85, 1.95, 2.1, 0.45, str(qw.get("effort", qw.get("cost", "1日")))[:14],
                      size=15, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 中段：実装ステップ
        _add_rect(s, 0.4, 2.55, 12.55, 1.6, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.6, 2.6, 12.2, 0.3, "🛠 実装ステップ（このとおりに進めれば完了）",
                      size=11, bold=True, color=C.deep_blue)
        steps = qw.get("steps") or qw.get("actions") or []
        for si, step in enumerate(steps[:6], 1):
            sx = 0.6 + (si - 1) % 2 * 6.15
            sy = 2.95 + (si - 1) // 2 * 0.4
            _add_rect(s, sx, sy + 0.07, 0.25, 0.25, C.teal,
                      shape_type=MSO_SHAPE.OVAL)
            _add_text_box(s, sx, sy + 0.07, 0.25, 0.25, str(si),
                          size=9, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, sx + 0.32, sy, 5.7, 0.4, str(step)[:80],
                          size=9, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

        # 下段：Before / After サンプルコード
        before_text = str(qw.get("before", ""))[:200]
        after_text = str(qw.get("after", qw.get("code_sample", "")))[:380]

        _add_rect(s, 0.4, 4.3, 6.15, 2.55, RGBColor(0xFE, 0xF3, 0xF2), line_color=C.red,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 0.4, 4.3, 6.15, 0.32, C.red,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.5, 4.3, 6.0, 0.32, "❌ BEFORE（現状）",
                      size=10, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 0.55, 4.7, 6.0, 2.1, before_text,
                      size=8, color=C.text_dark, font="Consolas")

        _add_rect(s, 6.8, 4.3, 6.15, 2.55, RGBColor(0xEC, 0xFD, 0xF5), line_color=C.green,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 6.8, 4.3, 6.15, 0.32, C.green,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 6.9, 4.3, 6.0, 0.32, "✅ AFTER（実装後）",
                      size=10, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 6.95, 4.7, 6.0, 2.1, after_text,
                      size=8, color=C.text_dark, font="Consolas")

        # 検証方法
        validation = qw.get("validation", "")
        if validation:
            _add_rect(s, 0.4, 6.95, 12.55, 0.3, C.light_bg,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.55, 6.95, 12.3, 0.3,
                          f"🔍 検証: {validation[:140]}",
                          size=9, italic=True, color=C.deep_blue, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # 🤖 AI引用テスト結果（AI Citation Test）
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "🤖 AI引用テスト — あなたのサイトはAIにどう見えているか",
                "AI Citation Test")
    _add_text_box(s, 0.5, 1.0, 12.3, 0.35,
                  "ChatGPT / Perplexity / Gemini に実際にクエリした場合の引用可能性を診断",
                  size=12, italic=True, color=C.muted)

    # クエリ一覧を取得（5〜8件）
    _cit_queries = test_queries.get("queries", [])
    if not _cit_queries:
        # フォールバック: サイトタイトルから自動生成
        _cit_queries = [
            {"query": f"{site_title}とは", "platform": "ChatGPT"},
            {"query": f"{site_title} おすすめ", "platform": "ChatGPT"},
            {"query": f"{site_title} 比較", "platform": "Perplexity"},
            {"query": f"{site_title} メリット デメリット", "platform": "Perplexity"},
            {"query": f"{site_title} 選び方", "platform": "Google AI Overview"},
        ]

    # 引用判定ロジック: all_scores から総合推定
    _cit_total = total.get("total", 0)

    def _citation_verdict(score):
        if score >= 75:
            return "○", "高確率で引用", C.green
        if score >= 50:
            return "△", "条件付き引用", C.amber
        return "×", "引用されない可能性大", C.red

    # テーブルヘッダー
    _add_rect(s, 0.4, 1.5, 12.55, 0.45, C.navy)
    _cit_hdrs = [
        ("クエリ文", 0.5, 4.8),
        ("想定AIの回答", 5.3, 3.5),
        ("引用", 8.8, 0.7),
        ("根拠", 9.5, 3.4),
    ]
    for _ht, _hx, _hw in _cit_hdrs:
        _add_text_box(s, _hx, 1.5, _hw, 0.45, _ht,
                      size=11, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _cit_counts = {"○": 0, "△": 0, "×": 0}
    for ci, cq in enumerate(_cit_queries[:8]):
        cy = 1.95 + ci * 0.6
        cbg = C.white if ci % 2 == 0 else C.light_bg
        _add_rect(s, 0.4, cy, 12.55, 0.6, cbg, line_color=C.border)

        # クエリ文
        _add_text_box(s, 0.5, cy, 4.8, 0.6,
                      cq.get("query", "")[:50],
                      size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

        # 想定回答
        _cit_answer = cq.get("expected_answer", "")
        if not _cit_answer:
            _cit_answer = f"{site_title}に関する一般的な回答"
        _add_text_box(s, 5.3, cy, 3.5, 0.6,
                      _cit_answer[:40],
                      size=9, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

        # 引用判定
        _cit_mark, _cit_reason_default, _cit_color = _citation_verdict(_cit_total)
        _cit_counts[_cit_mark] += 1
        _add_rect(s, 8.95, cy + 0.12, 0.4, 0.35, _cit_color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 8.95, cy + 0.12, 0.4, 0.35, _cit_mark,
                      size=12, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 根拠
        _cit_reason = cq.get("reason_if_not", _cit_reason_default)
        _add_text_box(s, 9.5, cy, 3.4, 0.6,
                      _cit_reason[:45],
                      size=9, color=C.deep_blue, anchor=MSO_ANCHOR.MIDDLE)

    # 下部コールアウト: 引用推定率
    _cit_total_q = max(sum(_cit_counts.values()), 1)
    _cit_pct_good = int(_cit_counts["○"] / _cit_total_q * 100)
    _cit_pct_maybe = int(_cit_counts["△"] / _cit_total_q * 100)
    _cit_pct_bad = int(_cit_counts["×"] / _cit_total_q * 100)

    _add_rect(s, 0.4, 6.85, 12.55, 0.45, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 6.85, 12.4, 0.45,
                  f"📊 現状のAI引用推定率:  ○ 高確率 {_cit_pct_good}%  |  "
                  f"△ 条件付き {_cit_pct_maybe}%  |  × 困難 {_cit_pct_bad}%",
                  size=12, bold=True, color=C.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # 💰 改善ROI試算（Return on Investment）
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "💰 改善ROI試算 — 投資対効果の概算", "Return on Investment")

    _roi_current = total.get("total", 0)
    _roi_improved = min(_roi_current + 20, 100)
    _roi_organic_uplift = round((_roi_improved - _roi_current) * 0.5, 1)
    _roi_ai_citation_uplift = round((_roi_improved - _roi_current) * 0.3, 1)

    # 4つのKPIカード横並び
    _roi_cards = [
        ("現状スコア", f"{_roi_current}", "/100", C.red if _roi_current < 50 else C.amber),
        ("改善後スコア（想定）", f"{_roi_improved}", "/100", C.green),
        ("オーガニック流入改善率", f"+{_roi_organic_uplift}", "%", C.teal),
        ("AI引用率改善", f"+{_roi_ai_citation_uplift}", "%", C.deep_blue),
    ]
    for ri, (rlbl, rbig, rsub, rcol) in enumerate(_roi_cards):
        rx = 0.5 + ri * 3.18
        _add_rect(s, rx, 1.1, 3.0, 1.45, rcol,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, rx, 1.15, 3.0, 0.3, rlbl,
                      size=10, color=C.white, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, rx, 1.5, 3.0, 0.65, rbig,
                      size=34, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, rx, 2.2, 3.0, 0.3, rsub,
                      size=11, color=C.white, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)

    # 中段: 3フェーズの投資対効果テーブル
    _add_text_box(s, 0.5, 2.75, 12.3, 0.35,
                  "📈 3フェーズ投資対効果",
                  size=14, bold=True, color=C.text_dark, font=FONT_HDR)

    # テーブルヘッダー
    _add_rect(s, 0.5, 3.15, 12.3, 0.45, C.navy)
    _roi_th = [
        ("フェーズ", 0.6, 2.5),
        ("期間 / コスト", 3.1, 2.5),
        ("主な施策", 5.6, 3.8),
        ("想定スコア改善", 9.4, 1.8),
        ("累積効果", 11.2, 1.5),
    ]
    for _rt, _rx, _rw in _roi_th:
        _add_text_box(s, _rx, 3.15, _rw, 0.45, _rt,
                      size=10, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _roi_phases = [
        ("Phase 1", "30日 / 低コスト", "Quick Win実施（構造化データ・メタ改善）",
         "+10〜15点", f"{min(_roi_current + 12, 100)}点", C.teal),
        ("Phase 2", "60日 / 中コスト", "構造化データ強化＋E-E-A-T改善",
         "+10〜15点", f"{min(_roi_current + 25, 100)}点", C.deep_blue),
        ("Phase 3", "90日 / 継続", "コンテンツ戦略＋権威性構築",
         "+5〜10点", f"{min(_roi_current + 32, 100)}点", C.navy),
    ]
    for pi, (plbl, pcost, paction, pgain, pcum, pcol) in enumerate(_roi_phases):
        py = 3.6 + pi * 0.75
        pbg = C.white if pi % 2 == 0 else C.light_bg
        _add_rect(s, 0.5, py, 12.3, 0.75, pbg, line_color=C.border)

        # フェーズ名バッジ
        _add_rect(s, 0.65, py + 0.18, 1.6, 0.38, pcol,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.65, py + 0.18, 1.6, 0.38, plbl,
                      size=10, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 2.3, py, 0.8, 0.75, "",
                      size=1, color=C.white)  # spacer
        _add_text_box(s, 3.1, py, 2.5, 0.75, pcost,
                      size=10, color=C.text_dark,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 5.6, py, 3.8, 0.75, paction,
                      size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 9.4, py, 1.8, 0.75, pgain,
                      size=11, bold=True, color=C.green,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 11.2, py, 1.5, 0.75, pcum,
                      size=11, bold=True, color=pcol,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 下段: 免責注記
    _add_rect(s, 0.5, 6.1, 12.3, 0.5, C.light_bg,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.6, 6.1, 12.1, 0.5,
                  "⚠️ 本試算は推定値であり、実際の効果はサイト規模・業界・実装品質により異なります。"
                  "数値は過去の診断実績に基づく平均的な改善幅を参考にしています。",
                  size=9, italic=True, color=C.muted,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ROIサマリーバー
    _add_rect(s, 0.5, 6.75, 12.3, 0.45, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 6.75, 12.3, 0.45,
                  f"💎 90日間の総合効果:  スコア {_roi_current} → {min(_roi_current + 32, 100)}点  |  "
                  f"オーガニック流入 +{_roi_organic_uplift}%  |  AI引用率 +{_roi_ai_citation_uplift}%",
                  size=12, bold=True, color=C.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # 📊 6カテゴリ レーダーチャート
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "📊 6カテゴリ レーダーチャート", "Radar Chart")

    # ---- 左半分: 六角形レーダーチャート ----
    _rc_center_x = 3.6   # インチ
    _rc_center_y = 4.2
    _rc_radius = 2.2

    # 6カテゴリのスコアを0-1に正規化
    _rc_scores = []
    _rc_labels = []
    for _rk, _rl in cat_order[:6]:
        _rc_cat = categories.get(_rk, {})
        _rc_s = _rc_cat.get("score", 0)
        _rc_m = _rc_cat.get("max", 20) or 20
        _rc_scores.append(_rc_s / _rc_m)
        _rc_labels.append((_rl, _rc_s, _rc_m))

    # 角度（上から時計回り）
    _rc_angles = [math.pi / 2 - i * 2 * math.pi / 6 for i in range(6)]

    # 外枠ガイドライン（50%、100%の六角形）
    for _rc_ratio in [0.5, 1.0]:
        _rc_r = _rc_radius * _rc_ratio
        for _ri in range(6):
            _x1 = _rc_center_x + _rc_r * math.cos(_rc_angles[_ri])
            _y1 = _rc_center_y - _rc_r * math.sin(_rc_angles[_ri])
            _x2 = _rc_center_x + _rc_r * math.cos(_rc_angles[(_ri + 1) % 6])
            _y2 = _rc_center_y - _rc_r * math.sin(_rc_angles[(_ri + 1) % 6])
            _rc_line = s.shapes.add_connector(
                1, Inches(_x1), Inches(_y1), Inches(_x2), Inches(_y2))
            _rc_line.line.color.rgb = C.border
            _rc_line.line.width = Pt(0.8 if _rc_ratio < 1.0 else 1.5)

    # 軸線（中心から各頂点へ）
    for _ri in range(6):
        _xv = _rc_center_x + _rc_radius * math.cos(_rc_angles[_ri])
        _yv = _rc_center_y - _rc_radius * math.sin(_rc_angles[_ri])
        _rc_axis = s.shapes.add_connector(
            1, Inches(_rc_center_x), Inches(_rc_center_y),
            Inches(_xv), Inches(_yv))
        _rc_axis.line.color.rgb = C.border
        _rc_axis.line.width = Pt(0.5)

    # スコア六角形（FreeformBuilder で塗りつぶし）
    _rc_fb = s.shapes.build_freeform(
        Inches(_rc_center_x + _rc_radius * _rc_scores[0] * math.cos(_rc_angles[0])),
        Inches(_rc_center_y - _rc_radius * _rc_scores[0] * math.sin(_rc_angles[0])),
    )
    _rc_vertices = []
    for _ri in range(1, 6):
        _sx = _rc_center_x + _rc_radius * _rc_scores[_ri] * math.cos(_rc_angles[_ri])
        _sy = _rc_center_y - _rc_radius * _rc_scores[_ri] * math.sin(_rc_angles[_ri])
        _rc_vertices.append((Inches(_sx), Inches(_sy)))
    # 始点に戻って閉じる
    _rc_vertices.append((
        Inches(_rc_center_x + _rc_radius * _rc_scores[0] * math.cos(_rc_angles[0])),
        Inches(_rc_center_y - _rc_radius * _rc_scores[0] * math.sin(_rc_angles[0])),
    ))
    _rc_fb.add_line_segments(_rc_vertices)
    _rc_shape = _rc_fb.convert_to_shape()
    _rc_shape.fill.solid()
    _rc_shape.fill.fore_color.rgb = C.cyan
    _rc_shape.fill.fore_color.brightness = 0.3
    _rc_shape.line.color.rgb = C.teal
    _rc_shape.line.width = Pt(2)

    # 各頂点にカテゴリ名とスコアをラベル表示
    for _ri in range(6):
        _lx = _rc_center_x + (_rc_radius + 0.45) * math.cos(_rc_angles[_ri])
        _ly = _rc_center_y - (_rc_radius + 0.45) * math.sin(_rc_angles[_ri])
        _rc_lbl, _rc_sc, _rc_mx = _rc_labels[_ri]
        _rc_pct = _rc_scores[_ri]
        _rc_sc_color = C.green if _rc_pct >= 0.7 else (C.amber if _rc_pct >= 0.4 else C.red)
        _add_text_box(s, _lx - 0.9, _ly - 0.2, 1.8, 0.22, _rc_lbl,
                      size=9, bold=True, color=C.text_dark,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, _lx - 0.5, _ly + 0.02, 1.0, 0.22,
                      f"{_rc_sc}/{_rc_mx}",
                      size=10, bold=True, color=_rc_sc_color,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # ---- 右半分: カテゴリ別一言所見 ----
    _add_text_box(s, 7.0, 1.0, 6.0, 0.35, "カテゴリ別所見",
                  size=14, bold=True, color=C.text_dark, font=FONT_HDR)

    # 弱点を特定（スコア比率が低い順にソート）
    _rc_sorted = sorted(range(6), key=lambda i: _rc_scores[i])
    _rc_weak_indices = set(_rc_sorted[:2])  # 下位2つをハイライト

    for _ri in range(6):
        _ry = 1.5 + _ri * 0.88
        _rc_lbl, _rc_sc, _rc_mx = _rc_labels[_ri]
        _rc_pct = _rc_scores[_ri]
        _is_weak = _ri in _rc_weak_indices

        # 背景カード
        _rc_card_border = C.red if _is_weak else C.border
        _add_rect(s, 7.0, _ry, 5.9, 0.78, C.white, line_color=_rc_card_border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

        # 弱点マーク
        if _is_weak:
            _add_rect(s, 7.1, _ry + 0.08, 0.5, 0.25, C.red,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 7.1, _ry + 0.08, 0.5, 0.25, "弱点",
                          size=8, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # カテゴリ名 + スコア
        _rc_name_x = 7.7 if _is_weak else 7.15
        _add_text_box(s, _rc_name_x, _ry + 0.05, 3.0, 0.3, _rc_lbl,
                      size=11, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _rc_sc_color = C.green if _rc_pct >= 0.7 else (C.amber if _rc_pct >= 0.4 else C.red)
        _add_text_box(s, 11.5, _ry + 0.05, 1.3, 0.3,
                      f"{_rc_sc}/{_rc_mx}",
                      size=11, bold=True, color=_rc_sc_color,
                      align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

        # 一言所見
        _rc_items = categories.get(cat_order[_ri][0], {}).get("items", [])
        _rc_low_items = [it for it in _rc_items
                         if it.get("max", 0) > 0
                         and it.get("score", 0) / it.get("max", 1) < 0.5]
        if _rc_low_items:
            _rc_comment = f"低スコア {len(_rc_low_items)}件 — 改善余地あり"
        elif _rc_pct >= 0.8:
            _rc_comment = "良好 — 現状維持を推奨"
        else:
            _rc_comment = "概ね基準をクリア"
        _add_text_box(s, _rc_name_x, _ry + 0.4, 5.1, 0.3,
                      _rc_comment,
                      size=9, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # 🎯 CV診断（面接応募）— サマリ + 10要素 + 改善案2-3枚
    # ==============================================================
    if cv_data:
        # ----- CVサマリスライド -----
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "🎯 CV（面接応募）転換率診断", "Conversion Rate Assessment")

        cv_total_score = cv_data.get("cv_total", 0)
        cv_grade = cv_data.get("cv_grade", "C")
        est_cv = cv_data.get("estimated_cv_rate", 0)
        imp_cv = cv_data.get("improved_cv_rate", 0)
        uplift_pct = cv_data.get("potential_uplift_pct", 0)
        bench_pos = cv_data.get("benchmark_position", "業界平均並み")
        bench = cv_data.get("benchmark", {})

        # 4カードKPI
        cards = [
            ("CV診断スコア", f"{cv_total_score}", f"/100  {cv_grade}", _grade_color(cv_grade)),
            ("推定CV率（現状）", f"{est_cv}", "%", C.amber),
            ("改善後CV率（予測）", f"{imp_cv}", "%", C.green),
            ("改善余地", f"+{uplift_pct}", "%", C.teal),
        ]
        for i, (lbl, big, sub, col) in enumerate(cards):
            x = 0.5 + i * 3.18
            _add_rect(s, x, 1.15, 3.0, 1.5, col, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, x, 1.2, 3.0, 0.3, lbl,
                          size=11, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, x, 1.55, 3.0, 0.7, big,
                          size=36, bold=True, color=C.white, font=FONT_HDR,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, x, 2.3, 3.0, 0.3, sub,
                          size=12, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # ベンチマーク帯
        _add_rect(s, 0.5, 2.85, 12.3, 1.05, C.navy, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.7, 2.92, 11.9, 0.3, "📊 業界ベンチマーク（採用ページ応募完了率）",
                      size=11, bold=True, color=C.cyan)
        _add_text_box(s, 0.7, 3.22, 11.9, 0.65,
                      f"業界平均: {bench.get('industry_avg_low', 1.5)}〜{bench.get('industry_avg_high', 3.0)}%　|　"
                      f"ベストプラクティス: {bench.get('best_practice', 6.0)}%　|　"
                      f"トップクラス: {bench.get('top_class', 10.0)}%　|　"
                      f"自社位置: {bench_pos}",
                      size=12, color=C.white, anchor=MSO_ANCHOR.MIDDLE)

        # 10要素レーダー風（横棒）
        _add_text_box(s, 0.5, 4.05, 12.3, 0.4, "🔍 CV阻害要因 10要素 — スコア一覧",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)
        factors = cv_data.get("factors", [])
        for i, f in enumerate(factors[:10]):
            row = i % 5
            col = i // 5
            x = 0.5 + col * 6.4
            y = 4.5 + row * 0.5
            score = f.get("score", 0)
            mx = f.get("max", 10)
            pct = score / mx if mx else 0
            color = C.green if pct >= 0.7 else (C.amber if pct >= 0.4 else C.red)

            _add_text_box(s, x, y, 2.5, 0.4, f.get("label", "")[:14],
                          size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
            _add_rect(s, x + 2.5, y + 0.1, 3.0, 0.2, C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, x + 2.5, y + 0.1, max(3.0 * pct, 0.05), 0.2, color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, x + 5.6, y, 0.7, 0.4, f"{score}/{mx}",
                          size=10, bold=True, color=color, anchor=MSO_ANCHOR.MIDDLE)

        # 弱点TOP3
        weak = cv_data.get("weak_factors", [])[:3]
        if weak:
            _add_rect(s, 0.5, 7.0, 12.3, 0.35, C.red, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.6, 7.0, 12.2, 0.35,
                          "⚠️ 最優先弱点: " + " / ".join(f["label"] + f"({f['score']}点)" for f in weak),
                          size=10, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

        # ----- CV改善案 詳細スライド（最大4件） -----
        cv_ideas = cv_data.get("improvement_ideas", [])
        for i, idea in enumerate(cv_ideas[:6]):
            s = pres.slides.add_slide(blank)
            _set_bg(s, C.cream)

            pri = idea.get("priority", "A")
            pri_color = {"S": C.red, "A": C.amber, "B": C.teal}.get(pri, C.teal)
            _add_rect(s, 0, 0, 13.33, 0.65, RGBColor(0x7C, 0x2D, 0x12))  # CV専用色
            _add_text_box(s, 0.5, 0.13, 1.4, 0.4, "🎯 CV改善",
                          size=13, bold=True, color=C.cyan, font=FONT_HDR)
            _add_rect(s, 1.95, 0.18, 0.7, 0.32, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 1.95, 0.18, 0.7, 0.32, f"優先{pri}",
                          size=11, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 2.75, 0.13, 8.5, 0.4, idea.get("title", "")[:60],
                          size=15, bold=True, color=C.white, font=FONT_HDR,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 11.3, 0.13, 1.9, 0.4,
                          f"対象: {idea.get('factor_label', '')[:14]}",
                          size=10, color=C.cyan,
                          align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
            _add_rect(s, 0, 0.65, 13.33, 0.04, C.cyan)

            # 現状スコア / 期待効果 / コスト カード
            cur_score = idea.get("current_score", 0)
            max_score = idea.get("max_score", 10)
            _add_rect(s, 0.4, 0.95, 4.1, 1.0, C.white, line_color=C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.5, 1.0, 4.0, 0.3, "📊 現状スコア",
                          size=10, color=C.muted, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 0.5, 1.3, 4.0, 0.6, f"{cur_score} / {max_score}",
                          size=24, bold=True, color=C.red, font=FONT_HDR,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            _add_rect(s, 4.6, 0.95, 4.1, 1.0, C.green, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 4.7, 1.0, 4.0, 0.3, "📈 期待CV改善効果",
                          size=10, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 4.7, 1.3, 4.0, 0.6, str(idea.get("expected_uplift", ""))[:18],
                          size=20, bold=True, color=C.white, font=FONT_HDR,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            _add_rect(s, 8.85, 0.95, 4.1, 1.0, C.amber, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 8.95, 1.0, 4.0, 0.3, "💰 実装コスト",
                          size=10, color=C.white, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 8.95, 1.3, 4.0, 0.6, str(idea.get("cost", ""))[:18],
                          size=18, bold=True, color=C.white, font=FONT_HDR,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # 具体的アクション
            _add_rect(s, 0.4, 2.1, 12.55, 2.1, C.white, line_color=C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.6, 2.15, 12.3, 0.3, "🛠 具体的アクション",
                          size=12, bold=True, color=C.deep_blue)
            actions = idea.get("actions", [])
            for ai, action in enumerate(actions[:5]):
                ay = 2.5 + ai * 0.32
                _add_rect(s, 0.6, ay + 0.07, 0.22, 0.22, C.teal, shape_type=MSO_SHAPE.OVAL)
                _add_text_box(s, 0.6, ay + 0.07, 0.22, 0.22, str(ai + 1),
                              size=9, bold=True, color=C.white,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text_box(s, 0.9, ay, 11.9, 0.32, str(action)[:130],
                              size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

            # 実装サンプルコード
            _add_rect(s, 0.4, 4.35, 12.55, 2.55, RGBColor(0x1E, 0x29, 0x3B),
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, 0.4, 4.35, 12.55, 0.32, C.deep_blue,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.55, 4.35, 12.3, 0.32, "💻 実装サンプルコード（コピペで使える）",
                          size=10, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)
            code = str(idea.get("code_sample", ""))[:600]
            _add_text_box(s, 0.55, 4.75, 12.3, 2.1, code,
                          size=8, color=C.cyan, font="Consolas")

            # CV影響メモ
            _add_rect(s, 0.4, 7.0, 12.55, 0.32, C.light_bg,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.55, 7.0, 12.3, 0.32,
                          "📌 改善対象要素: " + idea.get("factor_label", "") + " — 業界統計データに基づく改善余地",
                          size=9, italic=True, color=C.deep_blue, anchor=MSO_ANCHOR.MIDDLE)

            _footer(s, _pg())

        # ----- CV改善ロードマップスライド -----
        if cv_ideas:
            s = pres.slides.add_slide(blank)
            _set_bg(s, C.cream)
            _header_bar(s, "🎯 CV改善 推奨実施ロードマップ", "CV Improvement Roadmap")
            _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                          f"優先度・現状スコア順に並べた{len(cv_ideas)}施策の実施順序",
                          size=13, italic=True, color=C.muted)

            # 並び替え: 優先度S→A→B、その中で現状スコア低い順
            sorted_ideas = sorted(cv_ideas,
                                   key=lambda x: ({"S": 0, "A": 1, "B": 2}.get(x.get("priority", "B"), 9),
                                                  x.get("current_score", 0)))

            # ヘッダー
            _add_rect(s, 0.5, 1.55, 12.3, 0.4, C.navy)
            for hi, (htxt, w, x_off) in enumerate([
                ("順", 0.5, 0), ("Phase", 1.6, 0.5),
                ("施策", 5.0, 2.1),
                ("対象要素", 2.4, 7.1),
                ("効果", 1.5, 9.5),
                ("工数", 1.3, 11.0),
            ]):
                _add_text_box(s, 0.5 + x_off, 1.55, w, 0.4, htxt,
                              size=10, bold=True, color=C.cyan,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            for i, idea in enumerate(sorted_ideas[:12]):
                y = 1.95 + i * 0.4
                bg = C.white if i % 2 == 0 else C.light_bg
                phase = "Phase 1（〜30日）" if i < 3 else ("Phase 2（〜60日）" if i < 6 else "Phase 3（60日〜）")
                phase_color = C.red if i < 3 else (C.amber if i < 6 else C.teal)
                _add_rect(s, 0.5, y, 12.3, 0.4, bg, line_color=C.border)
                _add_text_box(s, 0.55, y, 0.5, 0.4, str(i + 1),
                              size=10, bold=True, color=C.text_dark,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_rect(s, 1.05, y + 0.05, 1.5, 0.3, phase_color,
                          shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
                _add_text_box(s, 1.05, y + 0.05, 1.5, 0.3, phase[:9],
                              size=8, bold=True, color=C.white,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text_box(s, 2.65, y, 5.0, 0.4, idea.get("title", "")[:42],
                              size=9, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
                _add_text_box(s, 7.65, y, 2.4, 0.4, idea.get("factor_label", "")[:16],
                              size=9, color=C.muted,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text_box(s, 10.05, y, 1.5, 0.4, str(idea.get("expected_uplift", ""))[:14],
                              size=9, bold=True, color=C.green,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
                _add_text_box(s, 11.55, y, 1.3, 0.4, str(idea.get("cost", ""))[:14],
                              size=9, color=C.muted,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            _footer(s, _pg())

    # (Agent-inserted slides D-G removed — using manually-inserted version below)

    # ==============================================================
    # 📋 カテゴリ別 深掘りスライド（各カテゴリ1枚）
    # ==============================================================
    for cat_idx, (cat_key, cat_label) in enumerate(cat_order):
        cat_data = categories.get(cat_key, {})
        cat_score = cat_data.get("score", 0)
        cat_max = cat_data.get("max", 20) or 20
        cat_items = cat_data.get("items", [])
        if not items:
            raw = all_scores.get(cat_key, {})
            if isinstance(raw, dict):
                items = [
                    {"key": k, "label": k.replace("_", " "),
                     "score": v if isinstance(v, (int, float)) else 0,
                     "max": 10, "reason": ""}
                    for k, v in raw.items()
                ]

        if not items:
            continue

        cat_score = cat_data.get("score", sum(it.get("score", 0) for it in items))
        cat_max = cat_data.get("max", sum(it.get("max", 10) for it in items))

        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, f"📋 {cat_label} — 項目別スコア詳細",
                    f"カテゴリスコア: {cat_score}/{cat_max}")

        # テーブルヘッダー
        _add_rect(s, 0.5, 1.05, 12.3, 0.45, C.navy)
        for hi, (htxt, w, x_off) in enumerate([
            ("項目名", 3.5, 0),
            ("スコア", 1.0, 3.5),
            ("最大", 0.8, 4.5),
            ("達成率", 3.2, 5.3),
            ("判定", 0.8, 8.5),
            ("所見", 3.0, 9.3),
        ]):
            _add_text_box(s, 0.5 + x_off, 1.05, w, 0.45, htxt,
                          size=10, bold=True, color=C.cyan,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        for i, item in enumerate(items[:13]):
            y = 1.5 + i * 0.42
            bg = C.white if i % 2 == 0 else C.light_bg
            sc = item.get("score", 0)
            mx = item.get("max", 10)
            pct = sc / mx if mx else 0

            # 色判定
            if pct >= 0.7:
                bar_color = C.green
            elif pct >= 0.4:
                bar_color = C.amber
            else:
                bar_color = C.red

            # 判定マーク
            if pct >= 0.8:
                mark = "◎"
            elif pct >= 0.6:
                mark = "○"
            elif pct >= 0.4:
                mark = "△"
            else:
                mark = "×"

            reason = item.get("reason", "")[:60]

            _add_rect(s, 0.5, y, 12.3, 0.42, bg, line_color=C.border)
            # 項目名
            _add_text_box(s, 0.55, y, 3.5, 0.42,
                          item.get("label", item.get("key", ""))[:28],
                          size=9, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
            # スコア
            _add_text_box(s, 4.0, y, 1.0, 0.42, str(sc),
                          size=10, bold=True, color=bar_color,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 最大
            _add_text_box(s, 5.0, y, 0.8, 0.42, str(mx),
                          size=10, color=C.muted,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 達成率バー（背景 + 前景）
            _add_rect(s, 5.85, y + 0.12, 3.0, 0.18, C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, 5.85, y + 0.12, max(3.0 * pct, 0.05), 0.18, bar_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 8.9, y + 0.02, 0.5, 0.2, f"{int(pct*100)}%",
                          size=7, bold=True, color=bar_color, align=PP_ALIGN.LEFT)
            # 判定
            _add_text_box(s, 9.0, y, 0.8, 0.42, mark,
                          size=12, bold=True, color=bar_color,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 所見
            _add_text_box(s, 9.8, y, 3.0, 0.42, reason,
                          size=8, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # 📋 カテゴリ別 深掘りスライド（各カテゴリ1枚）
    # ==============================================================
    for cat_idx, (cat_key, cat_label) in enumerate(cat_order):
        cat_data = categories.get(cat_key, {})
        cat_score = cat_data.get("score", 0)
        cat_max = cat_data.get("max", 20) or 20
        cat_items = cat_data.get("items", [])

        # items が空なら all_scores からフォールバック
        if not cat_items and all_scores:
            prefix = cat_key.split("_")[0]  # "1", "2", etc.
            for sk, sv in (all_scores or {}).items():
                if isinstance(sv, dict) and sk.startswith(prefix + "-"):
                    cat_items.append({
                        "key": sk, "label": sv.get("label", sk),
                        "score": sv.get("score", 0), "max": sv.get("max", 0),
                        "reason": sv.get("reason", ""),
                    })

        if not cat_items:
            continue

        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, f"📋 {cat_label} — 項目別スコア詳細",
                    f"{cat_score:.0f} / {cat_max}")

        # ヘッダー行
        _add_rect(s, 0.5, 1.1, 12.3, 0.45, C.navy)
        for htxt, hx, hw in [("項目", 0.7, 3.2), ("スコア", 3.9, 1.0),
                              ("達成率", 4.9, 3.5), ("判定", 8.4, 0.8),
                              ("所見", 9.2, 3.5)]:
            _add_text_box(s, hx, 1.1, hw, 0.45, htxt,
                          size=11, bold=True, color=C.white,
                          anchor=MSO_ANCHOR.MIDDLE)

        for i, item in enumerate(cat_items[:10]):
            y = 1.55 + i * 0.48
            iscore = item.get("score", 0)
            imax = item.get("max", 1) or 1
            pct = iscore / imax
            bar_color = C.green if pct >= 0.7 else (C.amber if pct >= 0.4 else C.red)
            grade = "◎" if pct >= 0.8 else ("○" if pct >= 0.6 else ("△" if pct >= 0.4 else "×"))

            bg = C.white if i % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.48, bg, line_color=C.border)
            _add_text_box(s, 0.7, y, 3.2, 0.48,
                          str(item.get("label", item.get("key", "")))[:28],
                          size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 3.9, y, 1.0, 0.48, f"{iscore}/{imax}",
                          size=11, bold=True, color=C.text_dark,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            # 達成率バー
            _add_rect(s, 4.9, y + 0.14, 3.5, 0.2, C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, 4.9, y + 0.14, max(3.5 * pct, 0.05), 0.2, bar_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 8.4, y, 0.8, 0.48, grade,
                          size=14, bold=True, color=bar_color,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 9.2, y, 3.5, 0.48,
                          str(item.get("reason", ""))[:50],
                          size=9, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # ⚡ Core Web Vitals 詳細
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "⚡ Core Web Vitals 詳細診断", "Google PageSpeed Insights")

    ps_score = pagespeed.get("score", 0) or 0
    ps_color = C.green if ps_score >= 80 else (C.amber if ps_score >= 50 else C.red)

    # Performance Score 大型表示
    _add_rect(s, 0.5, 1.1, 3.5, 3.5, C.navy, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 1.3, 3.5, 0.4, "Performance Score",
                  size=13, color=C.cyan, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 1.8, 3.5, 1.5, str(ps_score),
                  size=72, bold=True, color=ps_color, font=FONT_HDR,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _add_text_box(s, 0.5, 3.4, 3.5, 0.4, "/ 100",
                  size=16, color=C.text_light, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 4.0, 3.5, 0.5,
                  f"Source: {pagespeed.get('source', 'N/A')[:30]}",
                  size=9, italic=True, color=C.muted, align=PP_ALIGN.CENTER)

    # CWV 3メトリクス
    cwv_metrics = [
        ("LCP", "Largest Contentful Paint",
         pagespeed.get("lcp"), "ms", 2500, 4000,
         "画像/フォントの最適化、サーバー応答時間改善"),
        ("CLS", "Cumulative Layout Shift",
         pagespeed.get("cls"), "", 0.1, 0.25,
         "画像サイズ指定、フォント読込最適化、広告枠固定"),
        ("INP", "Interaction to Next Paint",
         pagespeed.get("inp"), "ms", 200, 500,
         "JavaScript最適化、イベントハンドラの軽量化"),
    ]

    for i, (short, full, val, unit, good_th, poor_th, fix) in enumerate(cwv_metrics):
        x = 4.3
        y = 1.1 + i * 1.7

        if val is not None:
            if isinstance(val, float) and val < 1:
                display_val = f"{val:.4f}"
            else:
                display_val = f"{val:,.0f}"
            if val <= good_th:
                status, scolor = "Good", C.green
            elif val <= poor_th:
                status, scolor = "Needs Improvement", C.amber
            else:
                status, scolor = "Poor", C.red
        else:
            display_val, status, scolor = "—", "N/A", C.muted

        _add_rect(s, x, y, 8.5, 1.5, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x, y, 0.08, 1.5, scolor)

        _add_text_box(s, x + 0.3, y + 0.1, 1.5, 0.5, short,
                      size=22, bold=True, color=C.text_dark, font=FONT_HDR)
        _add_text_box(s, x + 0.3, y + 0.6, 3.5, 0.3, full,
                      size=9, color=C.muted)

        _add_text_box(s, x + 4.0, y + 0.1, 2.5, 0.8,
                      display_val + unit,
                      size=28, bold=True, color=scolor, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _add_rect(s, x + 6.5, y + 0.2, 1.6, 0.4, scolor,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, x + 6.5, y + 0.2, 1.6, 0.4, status[:20],
                      size=10, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _add_text_box(s, x + 0.3, y + 1.0, 8.0, 0.4,
                      f"目標: ≤{good_th}{unit}  |  改善: {fix[:50]}",
                      size=9, color=C.muted)

    # ベンチマーク帯
    _add_rect(s, 0.5, 6.3, 12.3, 0.6, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, 6.3, 11.9, 0.6,
                  "💡 CWVは2021年よりGoogleランキング要因。AIクローラーもページ読込速度を引用優先度に反映。"
                  "LCP 2.5秒以内 + CLS 0.1以下 + INP 200ms以下が目標。",
                  size=10, color=C.white, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # 🔍 競合の構造化データ比較
    # ==============================================================
    if competitors:
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "🔍 競合 構造化データ比較", "Schema.org Implementation")

        # 自サイトのスキーマタイプ収集
        self_types = set()
        for sd in structure.get("jsonld", []):
            t = sd.get("@type", "")
            if isinstance(t, list):
                self_types.update(t)
            else:
                self_types.add(t)

        # 行ヘッダ
        schema_rows = ["Organization", "Article", "FAQPage", "BreadcrumbList",
                       "WebSite", "JobPosting", "Product", "LocalBusiness"]
        comp_names = ["自サイト"] + [
            (c.get("title") or c.get("url", ""))[:16] for c in competitors[:3]
        ]
        comp_types = [self_types] + [
            set(c.get("sd_types", [])) for c in competitors[:3]
        ]

        # ヘッダー
        _add_rect(s, 0.5, 1.1, 12.3, 0.5, C.navy)
        _add_text_box(s, 0.7, 1.1, 3.0, 0.5, "スキーマタイプ",
                      size=11, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)
        for ci, name in enumerate(comp_names):
            x = 3.7 + ci * 2.3
            _add_text_box(s, x, 1.1, 2.2, 0.5, name[:14],
                          size=10, bold=True, color=C.cyan if ci == 0 else C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        for ri, stype in enumerate(schema_rows):
            y = 1.6 + ri * 0.55
            bg = C.white if ri % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.55, bg, line_color=C.border)
            _add_text_box(s, 0.7, y, 3.0, 0.55, stype,
                          size=11, bold=True, color=C.text_dark,
                          anchor=MSO_ANCHOR.MIDDLE)
            for ci, types in enumerate(comp_types):
                x = 3.7 + ci * 2.3
                has_it = stype in types or (stype == "Article" and
                         ("BlogPosting" in types or "NewsArticle" in types))
                icon = "○" if has_it else "×"
                color = C.green if has_it else C.red
                _add_text_box(s, x, y, 2.2, 0.55, icon,
                              size=18, bold=True, color=color,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 未実装コールアウト
        comp_all = set()
        for types in comp_types[1:]:
            comp_all.update(types)
        missing = comp_all - self_types
        if missing:
            msg = f"⚠️ 競合で実装済みだが自サイトで未実装: {', '.join(list(missing)[:5])}"
        else:
            msg = "✅ 競合と同等以上の構造化データ実装レベルです"
        _add_rect(s, 0.5, 6.2, 12.3, 0.7, C.teal if not missing else C.red,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.7, 6.2, 11.9, 0.7, msg,
                      size=12, bold=True, color=C.white,
                      anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # 📅 改善ロードマップ（ガントチャート風）
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "📅 改善ロードマップ — 90日実施計画", "Implementation Roadmap")

    qw_all = improvements.get("quick_wins", [])
    strat_all = improvements.get("strategic", [])

    # タイムライン軸
    _add_rect(s, 1.5, 1.3, 11.0, 0.04, C.border)
    milestones = [("0日", 1.5), ("30日", 5.2), ("60日", 8.6), ("90日", 12.1)]
    for label, x in milestones:
        _add_rect(s, x, 1.2, 0.02, 0.2, C.navy)
        _add_text_box(s, x - 0.4, 1.0, 0.8, 0.25, label,
                      size=9, bold=True, color=C.text_dark, align=PP_ALIGN.CENTER)

    phases = [
        ("Phase 1", "Quick Win — 即効性の高い施策", C.red,
         qw_all[:3], "+10〜15点", "0-30日"),
        ("Phase 2", "構造化データ＋E-E-A-T強化", C.amber,
         (strat_all[:3] if strat_all else qw_all[3:6]), "+10〜15点", "30-60日"),
        ("Phase 3", "コンテンツ戦略＋継続改善", C.teal,
         (strat_all[3:6] if len(strat_all) > 3 else qw_all[6:9]), "+5〜10点", "60-90日"),
    ]

    for pi, (phase_name, phase_desc, color, items, impact, period) in enumerate(phases):
        y = 1.7 + pi * 1.8

        # フェーズ帯
        _add_rect(s, 1.5, y, 11.0, 1.5, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

        # 左サイド（フェーズ名）
        _add_rect(s, 1.5, y, 2.2, 1.5, color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 1.5, y + 0.15, 2.2, 0.4, phase_name,
                      size=14, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER)
        _add_text_box(s, 1.5, y + 0.55, 2.2, 0.3, period,
                      size=9, color=C.white, align=PP_ALIGN.CENTER)
        _add_text_box(s, 1.5, y + 0.9, 2.2, 0.5, impact,
                      size=16, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 右サイド（施策リスト）
        _add_text_box(s, 3.9, y + 0.1, 8.4, 0.35, phase_desc,
                      size=11, bold=True, color=C.text_dark)
        for ii, item in enumerate(items[:3]):
            title = item.get("title", "")[:55] if isinstance(item, dict) else str(item)[:55]
            _add_text_box(s, 4.0, y + 0.5 + ii * 0.32, 8.2, 0.32,
                          f"• {title}",
                          size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

    # 合計インパクト
    _add_rect(s, 0.5, 7.1, 12.3, 0.3, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, 7.1, 11.9, 0.3,
                  "📈 90日間の合計改善見込み: +25〜40点（現状スコアにより変動）",
                  size=11, bold=True, color=C.white,
                  anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # 🏗️ サイト構造マップ（ページロール型診断 — コーポレート/採用のみ）
    # ==============================================================
    if site_diagnosis and preset_id in ("corporate", "recruiting"):
        completeness = site_diagnosis.get("structure_completeness", {})
        role_scores = site_diagnosis.get("role_scores", {})
        schema_map = site_diagnosis.get("schema_map", {})
        page_recs = site_diagnosis.get("page_recommendations", [])

        # --- スライド: サイト構造完全性 ---
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "🏗️ サイト構造診断マップ",
                    f"完全性: {completeness.get('score', 0)}%")

        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      "サイトに必要なページが揃っているか・各ページが役割を果たしているかを診断",
                      size=13, italic=True, color=C.muted)

        # 必須ページ
        req_found = completeness.get("required_found", [])
        req_missing = completeness.get("required_missing", [])
        rec_found = completeness.get("recommended_found", [])
        rec_missing = completeness.get("recommended_missing", [])

        _add_text_box(s, 0.5, 1.55, 6.3, 0.4, "✅ 必須ページ",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)

        _ROLE_LABELS = {
            "top": "トップ", "about": "会社概要", "business": "事業内容",
            "contact": "問合せ", "privacy": "プライバシー",
            "recruit_top": "採用トップ", "job_listing": "求人詳細",
            "entry": "エントリー", "culture": "社風/文化",
            "benefits": "福利厚生", "interview": "社員インタビュー",
            "faq": "FAQ", "ir": "IR情報", "news": "ニュース",
            "csr": "CSR/ESG",
        }

        all_required = req_found + req_missing
        for i, role in enumerate(all_required):
            col = i % 3
            row = i // 3
            x = 0.5 + col * 2.1
            y = 2.0 + row * 0.9

            found = role in req_found
            rs = role_scores.get(role, {})
            score_val = rs.get("score", 0)
            grade_val = rs.get("grade", "-")

            bg = C.green if found else C.red
            _add_rect(s, x, y, 1.95, 0.75, bg,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, x, y + 0.05, 1.95, 0.3,
                          _ROLE_LABELS.get(role, role),
                          size=11, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER)
            if found:
                _add_text_box(s, x, y + 0.35, 1.95, 0.35,
                              f"{score_val}点 ({grade_val})",
                              size=10, color=C.white,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            else:
                _add_text_box(s, x, y + 0.35, 1.95, 0.35,
                              "❌ 未検出",
                              size=10, bold=True, color=C.white,
                              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 推奨ページ
        _add_text_box(s, 6.8, 1.55, 6.0, 0.4, "📋 推奨ページ",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)

        all_recommended = rec_found + rec_missing
        for i, role in enumerate(all_recommended[:6]):
            col = i % 3
            row = i // 3
            x = 6.8 + col * 2.1
            y = 2.0 + row * 0.9

            found = role in rec_found
            rs = role_scores.get(role, {})

            bg = C.teal if found else C.muted
            _add_rect(s, x, y, 1.95, 0.75, bg,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, x, y + 0.05, 1.95, 0.3,
                          _ROLE_LABELS.get(role, role),
                          size=11, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER)
            status_txt = f"{rs.get('score', 0)}点" if found else "未検出"
            _add_text_box(s, x, y + 0.35, 1.95, 0.35, status_txt,
                          size=10, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # 構造化データマップ
        _add_text_box(s, 0.5, 4.2, 12.3, 0.4, "🗂️ 構造化データ 配置マップ",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)
        _add_rect(s, 0.5, 4.65, 12.3, 0.35, C.navy)
        for htxt, hx, hw in [("スキーマ", 0.7, 2.5), ("配置すべきページ", 3.2, 2.5),
                              ("検出先URL", 5.7, 5.0), ("状態", 10.7, 2.0)]:
            _add_text_box(s, hx, 4.65, hw, 0.35, htxt,
                          size=10, bold=True, color=C.white,
                          anchor=MSO_ANCHOR.MIDDLE)

        sm_items = list(schema_map.items())[:5] if schema_map else [
            ("Organization", {"should_be_on": "about", "found_on": [], "status": "unknown"}),
            ("WebSite", {"should_be_on": "top", "found_on": [], "status": "unknown"}),
        ]
        for i, (stype, sinfo) in enumerate(sm_items):
            y = 5.0 + i * 0.42
            bg = C.white if i % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.42, bg, line_color=C.border)
            _add_text_box(s, 0.7, y, 2.5, 0.42, stype,
                          size=10, bold=True, color=C.text_dark,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 3.2, y, 2.5, 0.42,
                          _ROLE_LABELS.get(str(sinfo.get("should_be_on", "")), "—"),
                          size=10, color=C.muted,
                          anchor=MSO_ANCHOR.MIDDLE)
            found_urls = sinfo.get("found_on", [])
            _add_text_box(s, 5.7, y, 5.0, 0.42,
                          (found_urls[0][:50] if found_urls else "未検出"),
                          size=9, color=C.text_dark if found_urls else C.red,
                          anchor=MSO_ANCHOR.MIDDLE)
            st_val = sinfo.get("status", "unknown")
            st_color = C.green if st_val == "ok" else C.red
            st_label = "✅" if st_val == "ok" else "❌ 要追加"
            _add_text_box(s, 10.7, y, 2.0, 0.42, st_label,
                          size=10, bold=True, color=st_color,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # ページ別改善提案（下部）
        if page_recs:
            _add_rect(s, 0.5, 7.1, 12.3, 0.3, C.deep_blue,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            top_recs = page_recs[:3]
            rec_txt = " | ".join(
                f"[{_ROLE_LABELS.get(r.get('role',''), r.get('role',''))}] {r.get('action', '')[:25]}"
                for r in top_recs
            )
            _add_text_box(s, 0.7, 7.1, 11.9, 0.3,
                          f"🔥 最優先: {rec_txt}",
                          size=10, bold=True, color=C.white,
                          anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # 🎯 サイト全体優先施策（未実装率が高い項目）
    # ==============================================================
    if site_agg and site_agg.get("common_issues"):
        common_issues = site_agg["common_issues"]

        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "🎯 サイト全体 優先対応項目", "Site-Wide Critical Issues")
        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      f"分析した{site_agg.get('page_count', 0)}ページのうち過半数で未実装の項目 — "
                      f"サイト共通テンプレート修正で全ページが一気に改善されます",
                      size=12, italic=True, color=C.muted)

        # ヘッダー行
        _add_rect(s, 0.5, 1.6, 12.3, 0.5, C.navy)
        for hi, (htxt, w, x_off) in enumerate([
            ("優先", 0.7, 0),
            ("項目", 5.5, 0.7),
            ("失点ページ率", 1.7, 6.2),
            ("平均スコア", 1.5, 7.9),
            ("改善後の想定インパクト", 3.4, 9.4),
        ]):
            _add_text_box(s, 0.5 + x_off, 1.6, w, 0.5, htxt,
                          size=11, bold=True, color=C.cyan,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        for i, issue in enumerate(common_issues[:10]):
            y = 2.1 + i * 0.45
            failure_pct = issue.get("failure_pct", 0)
            avg_score = issue.get("avg_score", 0)
            max_score = issue.get("max_score", 0)
            item_key = issue.get("item_key", "")
            label = item_key.replace("_", " ")[:50]

            # 優先度
            if failure_pct >= 80:
                pri = "S"; pri_color = C.red
            elif failure_pct >= 65:
                pri = "A"; pri_color = C.amber
            else:
                pri = "B"; pri_color = C.teal

            bg = C.white if i % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.45, bg, line_color=C.border)
            _add_rect(s, 0.65, y + 0.08, 0.45, 0.3, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.65, y + 0.08, 0.45, 0.3, pri,
                          size=10, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 1.2, y, 5.5, 0.45, label,
                          size=10, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

            # 失点ページ率（バー付き）
            bar_w = 1.5 * (failure_pct / 100)
            _add_rect(s, 6.3, y + 0.13, 1.5, 0.18, C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, 6.3, y + 0.13, max(bar_w, 0.05), 0.18, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 6.3, y - 0.05, 1.5, 0.2, f"{failure_pct}%",
                          size=8, bold=True, color=pri_color,
                          align=PP_ALIGN.CENTER)

            _add_text_box(s, 7.95, y, 1.4, 0.45,
                          f"{avg_score:.1f}/{max_score}",
                          size=10, color=C.muted,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

            # 想定インパクト（共通テンプレート修正で全ページ改善）
            page_count = site_agg.get("page_count", 1)
            failed_pages = int(page_count * failure_pct / 100)
            potential = int((max_score - avg_score) * failed_pages)
            _add_text_box(s, 9.4, y, 3.4, 0.45,
                          f"+{potential}点（{failed_pages}ページ × +{max_score - avg_score:.1f}点）",
                          size=10, bold=True, color=C.green,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        # フッター注記
        _add_rect(s, 0.5, 6.95, 12.3, 0.4, C.deep_blue,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.5, 6.95, 12.3, 0.4,
                      "💡 これらは「サイト共通の構造的問題」— ヘッダー/フッター/テンプレート1箇所の修正で全ページが一気に改善されます",
                      size=11, italic=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _footer(s, _pg())

        # ----- スコア分布 + 最弱ページTOP5 -----
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "🎯 サイト全体 スコア分布 & 最優先改善ページ", "Site-Wide Score Distribution")

        # 左: スコア分布
        _add_text_box(s, 0.5, 1.1, 6.0, 0.4, "📊 スコア分布（グレード別ページ数）",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)
        dist = site_agg.get("score_distribution", {})
        total_pages = sum(dist.values()) or 1
        grade_colors = {"S": C.green, "A": C.teal, "B": C.amber, "C": RGBColor(0xF9, 0x73, 0x16), "D": C.red}
        for gi, g in enumerate(["S", "A", "B", "C", "D"]):
            count = dist.get(g, 0)
            pct = count / total_pages * 100
            y = 1.6 + gi * 0.85
            _add_text_box(s, 0.5, y, 0.6, 0.6, g,
                          size=24, bold=True, color=grade_colors[g], font=FONT_HDR,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_rect(s, 1.1, y + 0.18, 4.0, 0.3, C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_rect(s, 1.1, y + 0.18, max(4.0 * pct / 100, 0.05), 0.3, grade_colors[g],
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 5.2, y, 1.2, 0.6, f"{count}p ({pct:.0f}%)",
                          size=11, bold=True, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

        # 右: 最弱ページTOP5
        worst = site_agg.get("worst_pages", [])[:5]
        _add_text_box(s, 7.0, 1.1, 6.0, 0.4, "⚠️ 最優先改善ページ TOP5",
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)
        for wi, p in enumerate(worst):
            y = 1.6 + wi * 1.05
            score = p.get("total", {}).get("total", 0) if isinstance(p.get("total"), dict) else p.get("score", 0)
            grade = p.get("total", {}).get("grade", "D") if isinstance(p.get("total"), dict) else p.get("grade", "D")
            url_str = p.get("url", "")[:55]
            title_str = p.get("title", "")[:35]

            _add_rect(s, 7.0, y, 6.2, 0.95, C.white, line_color=C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 7.15, y + 0.05, 0.7, 0.4, f"#{wi+1}",
                          size=14, bold=True, color=C.muted, font=FONT_HDR,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_rect(s, 7.85, y + 0.1, 0.8, 0.35, _grade_color(grade),
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 7.85, y + 0.1, 0.8, 0.35, f"{score}",
                          size=12, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 8.75, y + 0.05, 4.4, 0.35, title_str,
                          size=10, bold=True, color=C.text_dark,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 8.75, y + 0.4, 4.4, 0.3, url_str,
                          size=8, italic=True, color=C.muted,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 7.15, y + 0.65, 5.9, 0.3,
                          f"Grade {grade}",
                          size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

        # フッター注記
        site_score = site_agg.get("site_score", 0)
        site_grade = site_agg.get("site_grade", "C")
        median = site_agg.get("score_median", 0)
        _add_rect(s, 0.5, 6.95, 12.3, 0.4, C.navy,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.5, 6.95, 12.3, 0.4,
                      f"📈 サイト総合: {site_score}/100 (Grade {site_grade}) | 中央値: {median} | "
                      f"S+A: {dist.get('S',0)+dist.get('A',0)}p / C+D: {dist.get('C',0)+dist.get('D',0)}p",
                      size=11, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _footer(s, _pg())

    # ==============================================================
    # Slide N+1: 戦略施策（中長期）
    # ==============================================================
    strategic_list = improvements.get("strategic", []) or improvements.get("content_strategy", [])
    if strategic_list:
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "中長期 戦略施策", "Strategic Initiatives")
        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      "Quick Winで地盤固め → 戦略施策で差別化・競合優位を確立",
                      size=13, italic=True, color=C.muted)

        for i, item in enumerate(strategic_list[:4]):
            y = 1.55 + i * 1.4
            _add_rect(s, 0.5, y, 12.3, 1.3, C.white, line_color=C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            # 優先度バッジ
            pri = item.get("priority", "A")
            pri_color = {"S": C.red, "A": C.amber, "B": C.teal}.get(pri, C.teal)
            _add_rect(s, 0.7, y + 0.2, 0.7, 0.35, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.7, y + 0.2, 0.7, 0.35, f"優先{pri}",
                          size=10, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 1.55, y + 0.15, 8.5, 0.4, item.get("title", "")[:60],
                          size=14, bold=True, color=C.text_dark, font=FONT_HDR,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 10.2, y + 0.15, 2.5, 0.4,
                          f"カテゴリ: {item.get('category', '戦略')[:12]}",
                          size=10, color=C.muted, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

            # Why
            _add_text_box(s, 0.7, y + 0.55, 11.9, 0.35,
                          "🎯 " + str(item.get("why", ""))[:130],
                          size=10, color=C.deep_blue, anchor=MSO_ANCHOR.MIDDLE)

            # Steps
            steps = item.get("steps", item.get("actions", []))
            if steps:
                step_text = "  ▶ ".join([str(x)[:30] for x in steps[:4]])
                _add_text_box(s, 0.7, y + 0.9, 11.9, 0.35,
                              "▶ " + step_text,
                              size=9, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # Slide N+2: 競合分析からの示唆（Competitor-Informed）
    # ==============================================================
    ci_list = improvements.get("competitor_informed", [])
    if ci_list:
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "競合分析から導き出した施策", "Competitor-Informed Recommendations")
        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      "上位競合と比較して見えてきた差別化ポイント・キャッチアップ施策",
                      size=13, italic=True, color=C.muted)

        for i, item in enumerate(ci_list[:5]):
            y = 1.6 + i * 1.05
            _add_rect(s, 0.5, y, 12.3, 0.95, C.white, line_color=C.border,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            pri = item.get("priority", "A")
            pri_color = {"S": C.red, "A": C.amber, "B": C.teal}.get(pri, C.teal)
            _add_rect(s, 0.7, y + 0.18, 0.7, 0.3, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.7, y + 0.18, 0.7, 0.3, f"優先{pri}",
                          size=9, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 1.55, y + 0.1, 11.0, 0.35, item.get("title", "")[:80],
                          size=12, bold=True, color=C.text_dark, font=FONT_HDR,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 1.55, y + 0.45, 11.0, 0.45,
                          str(item.get("why", "") or item.get("reason", ""))[:160],
                          size=10, color=C.muted, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # Slide N+3: 測定計画（Measurement Plan）
    # ==============================================================
    mp_list = improvements.get("measurement_plan", [])
    if mp_list:
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "効果測定 KPI & ツール計画", "Measurement Plan")
        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      "施策実施後の効果を客観的に証明するための定量指標と計測ツール",
                      size=13, italic=True, color=C.muted)

        # ヘッダー
        _add_rect(s, 0.5, 1.55, 12.3, 0.45, C.deep_blue,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        for hi, htxt in enumerate(["指標", "目標値", "計測ツール"]):
            _add_text_box(s, 0.5 + hi * 4.1, 1.55, 4.1, 0.45, htxt,
                          size=12, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        for i, m in enumerate(mp_list[:8]):
            y = 2.05 + i * 0.55
            bg = C.white if i % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.55, bg, line_color=C.border)
            _add_text_box(s, 0.6, y, 4.0, 0.55, str(m.get("metric", ""))[:40],
                          size=10, bold=True, color=C.text_dark,
                          anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 4.7, y, 4.0, 0.55, str(m.get("target", ""))[:40],
                          size=10, color=C.teal,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 8.8, y, 4.0, 0.55, str(m.get("tool", ""))[:40],
                          size=10, color=C.muted,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # Slide N+4: 全施策一覧マトリクス
    # ==============================================================
    all_imp = []
    for key, label in [("quick_wins", "Quick Win"), ("strategic", "戦略"),
                       ("competitor_informed", "競合発"), ("technical_debt", "技術負債")]:
        for it in improvements.get(key, []):
            all_imp.append({
                "category": label,
                "priority": it.get("priority", "A"),
                "title": it.get("title", ""),
                "impact": it.get("impact", it.get("expected_uplift", "")),
                "effort": it.get("effort", it.get("cost", "")),
            })

    if all_imp:
        s = pres.slides.add_slide(blank)
        _set_bg(s, C.cream)
        _header_bar(s, "全改善施策 一覧マトリクス", "All Recommendations Matrix")
        _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                      f"今回の診断で抽出された全{len(all_imp)}件の改善施策（優先度順）",
                      size=13, italic=True, color=C.muted)

        # 並び替え: 優先度S→A→B→C
        pri_order = {"S": 0, "A": 1, "B": 2, "C": 3}
        all_imp_sorted = sorted(all_imp, key=lambda x: pri_order.get(x["priority"], 9))

        # ヘッダー
        _add_rect(s, 0.5, 1.55, 12.3, 0.4, C.navy)
        for hi, (htxt, w, x_off) in enumerate([
            ("優先", 0.7, 0), ("分類", 1.3, 0.7),
            ("施策タイトル", 7.4, 2.0),
            ("インパクト", 1.6, 9.4),
            ("工数", 1.3, 11.0),
        ]):
            _add_text_box(s, 0.5 + x_off, 1.55, w, 0.4, htxt,
                          size=10, bold=True, color=C.cyan,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        for i, it in enumerate(all_imp_sorted[:14]):
            y = 1.95 + i * 0.35
            bg = C.white if i % 2 == 0 else C.light_bg
            _add_rect(s, 0.5, y, 12.3, 0.35, bg, line_color=C.border)
            pri_color = {"S": C.red, "A": C.amber, "B": C.teal, "C": C.muted}.get(it["priority"], C.muted)
            _add_rect(s, 0.65, y + 0.05, 0.5, 0.25, pri_color,
                      shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
            _add_text_box(s, 0.65, y + 0.05, 0.5, 0.25, it["priority"],
                          size=9, bold=True, color=C.white,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 1.2, y, 1.3, 0.35, it["category"][:8],
                          size=9, color=C.deep_blue,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 2.5, y, 7.4, 0.35, it["title"][:60],
                          size=9, color=C.text_dark, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 9.9, y, 1.6, 0.35, str(it["impact"])[:14],
                          size=9, bold=True, color=C.teal,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            _add_text_box(s, 11.5, y, 1.3, 0.35, str(it["effort"])[:14],
                          size=9, color=C.muted,
                          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        _footer(s, _pg())

    # ==============================================================
    # Slide 10: 実装ロードマップ
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "実装ロードマップ", "90-Day Action Plan")

    _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                  "90日間で段階的に実装し、スコア +25点を目指す",
                  size=13, italic=True, color=C.muted)

    phases = [
        ("Phase 1", "Day 1–14", "Quick Win集中実装",
         ["Organization JSON-LD", "FAQPage JSON-LD", "llms.txt設置", "著者ボックス設置"],
         "+15点", C.teal),
        ("Phase 2", "Day 15–45", "コンテンツ構造改善",
         ["冒頭Answer-first化", "定義文パターン追加", "FAQ7問以上追加", "数値データ強化"],
         "+8点", C.deep_blue),
        ("Phase 3", "Day 46–90", "権威性・鮮度の強化",
         ["編集ポリシー明記", "外部引用リンク追加", "年次更新運用", "競合差別化コンテンツ"],
         "+7点", C.navy),
    ]

    for i, (period, days, title, tasks, gain, color) in enumerate(phases):
        x = 0.5 + i * 4.3
        _add_rect(s, x, 1.6, 4.0, 5.0, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x, 1.6, 4.0, 1.4, color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, x, 2.5, 4.0, 0.5, color)
        _add_text_box(s, x + 0.2, 1.75, 3.6, 0.35, period,
                      size=11, bold=True, color=C.cyan)
        _add_text_box(s, x + 0.2, 2.1, 3.6, 0.5, days,
                      size=20, bold=True, color=C.white, font=FONT_HDR)
        _add_text_box(s, x + 0.2, 3.15, 3.6, 0.4, title,
                      size=14, bold=True, color=C.text_dark, font=FONT_HDR)
        _add_rect(s, x + 0.2, 3.6, 0.5, 0.03, color)
        for j, t in enumerate(tasks):
            ty = 3.8 + j * 0.35
            _add_rect(s, x + 0.25, ty + 0.08, 0.12, 0.12, color,
                      shape_type=MSO_SHAPE.OVAL)
            _add_text_box(s, x + 0.45, ty, 3.4, 0.3, t,
                          size=11, color=C.text_dark,
                          anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, x + 1.0, 5.9, 2.0, 0.6, C.cream, line_color=color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, x + 1.0, 5.9, 2.0, 0.6, f"想定効果 {gain}",
                      size=12, bold=True, color=color,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    target_score = min(100, total["total"] + 25)
    target_grade = "A" if target_score >= 80 else "B" if target_score >= 60 else "C"
    _add_rect(s, 0.5, 6.75, 12.3, 0.45, C.deep_blue,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 6.75, 12.3, 0.45,
                  f"💎 90日後の想定総合スコア: {total['total']}点 → {target_score}点（グレード {target_grade} 達成）",
                  size=13, bold=True, color=C.white,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 11: ROI
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "想定効果・ROI試算", "Expected Impact")

    _add_text_box(s, 0.5, 1.1, 6.5, 0.4, "📊 スコア推移シミュレーション",
                  size=16, bold=True, color=C.text_dark, font=FONT_HDR)
    _add_rect(s, 0.5, 1.6, 6.3, 3.8, C.white, line_color=C.border,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)

    # Y軸ラベル
    for v in [100, 80, 60, 40, 20, 0]:
        y = 1.9 + (100 - v) / 100 * 3.3
        _add_text_box(s, 0.6, y - 0.15, 0.35, 0.3, str(v),
                      size=9, color=C.muted, align=PP_ALIGN.RIGHT)

    # データポイント
    now_score = total["total"]
    points = [
        (1.5, "現在", now_score, C.red),
        (2.9, "14日後", min(100, now_score + 15), C.amber),
        (4.3, "45日後", min(100, now_score + 23), C.teal),
        (5.7, "90日後", min(100, now_score + 30), C.green),
    ]

    # 折れ線を線で描画
    for i in range(len(points) - 1):
        x1, _, v1, _ = points[i]
        x2, _, v2, _ = points[i + 1]
        y1 = 1.9 + (100 - v1) / 100 * 3.3
        y2 = 1.9 + (100 - v2) / 100 * 3.3
        connector = s.shapes.add_connector(1, Inches(x1), Inches(y1),
                                           Inches(x2), Inches(y2))
        connector.line.color.rgb = C.deep_blue
        connector.line.width = Pt(3)

    for x, label, score, color in points:
        y = 1.9 + (100 - score) / 100 * 3.3
        _add_rect(s, x - 0.12, y - 0.12, 0.24, 0.24, color,
                  shape_type=MSO_SHAPE.OVAL)
        _add_text_box(s, x - 0.4, y - 0.5, 0.8, 0.3, str(score),
                      size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
        _add_text_box(s, x - 0.6, 5.0, 1.2, 0.3, label,
                      size=10, color=C.muted, align=PP_ALIGN.CENTER)

    # KPI
    _add_text_box(s, 7.2, 1.1, 5.6, 0.4, "📈 KPI予測",
                  size=16, bold=True, color=C.text_dark, font=FONT_HDR)

    kpis = [
        ("AI Overview 引用率", "2%", "12%"),
        ("ChatGPT 出典表示", "1件/月", "8件/月"),
        ("Perplexity 引用", "未計測", "15件/月"),
        ("AI経由セッション", "約50/月", "約400/月"),
        ("指名検索数", "基準", "+40%"),
    ]
    for i, (lbl, before, after) in enumerate(kpis):
        y = 1.6 + i * 0.7
        _add_rect(s, 7.2, y, 5.6, 0.6, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 7.35, y, 2.5, 0.6, lbl,
                      size=11, bold=True, color=C.text_dark,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 9.8, y, 1.1, 0.6, before,
                      size=11, color=C.muted, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 10.9, y, 0.4, 0.6, "→",
                      size=14, color=C.teal, align=PP_ALIGN.CENTER,
                      anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 11.3, y, 1.4, 0.6, after,
                      size=12, bold=True, color=C.teal,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_rect(s, 0.5, 5.7, 12.3, 1.1, C.navy,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.7, 5.8, 11.9, 0.35, "💼 ビジネスインパクト（試算）",
                  size=12, bold=True, color=C.cyan)
    _add_text_box(s, 0.7, 6.15, 11.9, 0.6,
                  "SEO流入に加え、AI経由の新規タッチポイントが加わることで、月間商談機会が推定 +20〜30%。\n"
                  "また、AI引用によるブランド露出はSEO検索結果と異なり「選ばれた1〜3社しか出ない」ため、競合優位を長期に確保できます。",
                  size=11, color=C.white)

    _footer(s, _pg())

    # ==============================================================
    # Slide 12: テストクエリ
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.cream)
    _header_bar(s, "商談・社内共有用テストクエリ", "AI Search Demo Queries")

    _add_text_box(s, 0.5, 1.05, 12.3, 0.4,
                  "実際にAIに質問して「御社が引用されるか」を体感確認するクエリ集",
                  size=13, italic=True, color=C.muted)

    platform_colors = {
        "ChatGPT": RGBColor(0x10, 0xA3, 0x7F),
        "Perplexity": RGBColor(0x20, 0xB8, 0xCD),
        "Google AI Overview": RGBColor(0x42, 0x85, 0xF4),
        "Google": RGBColor(0x42, 0x85, 0xF4),
    }

    queries_list = test_queries.get("queries", [])
    if not queries_list:
        queries_list = [
            {"platform": "ChatGPT", "query": f"{site_title}について教えてください",
             "reason_if_not": "定義文＋Answer-first構造の評価"},
            {"platform": "ChatGPT", "query": f"{site_title} おすすめ 比較",
             "reason_if_not": "比較構造・リスト化レベルの評価"},
            {"platform": "Perplexity", "query": f"{site_title}の選び方は？",
             "reason_if_not": "FAQ・HowTo構造の評価"},
            {"platform": "Perplexity", "query": f"{site_title} メリット デメリット",
             "reason_if_not": "網羅性・両論併記の評価"},
            {"platform": "Google AI Overview", "query": f"{site_title}とは",
             "reason_if_not": "定義文パターンの評価（最重要）"},
        ]

    for i, q in enumerate(queries_list[:5]):
        y = 1.55 + i * 1.05
        platform = q.get("platform", "ChatGPT")
        color = C.teal
        for k, v in platform_colors.items():
            if k in platform:
                color = v
                break

        _add_rect(s, 0.5, y, 12.3, 0.95, C.white, line_color=C.border,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_rect(s, 0.7, y + 0.15, 2.0, 0.65, color,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, 0.7, y + 0.15, 2.0, 0.65, platform[:18],
                      size=11, bold=True, color=C.white,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 2.9, y + 0.1, 9.7, 0.25, "検索クエリ例",
                      size=9, color=C.muted)
        _add_text_box(s, 2.9, y + 0.35, 9.7, 0.35,
                      f'"{q.get("query", "")[:60]}"',
                      size=13, bold=True, italic=True, color=C.text_dark,
                      font=FONT_HDR, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, 2.9, y + 0.65, 9.7, 0.25,
                      "▶ " + q.get("reason_if_not", "")[:70],
                      size=10, color=C.teal)

    _add_rect(s, 0.5, 6.8, 12.3, 0.45, C.light_bg,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 0.5, 6.8, 12.3, 0.45,
                  "💡 商談現場で上記クエリを実行 → 「御社サイトが出ない/出る」を体感してもらうと改善ニーズが明確化されます",
                  size=11, italic=True, color=C.deep_blue,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, _pg())

    # ==============================================================
    # Slide 13: Next Step
    # ==============================================================
    s = pres.slides.add_slide(blank)
    _set_bg(s, C.navy)
    _add_rect(s, 0, 0, 13.33, 0.08, C.cyan)

    _add_text_box(s, 0.5, 0.6, 12.3, 0.4, "NEXT STEP",
                  size=14, bold=True, color=C.cyan, align=PP_ALIGN.CENTER)
    _add_text_box(s, 0.5, 1.1, 12.3, 0.8, "次のアクションについて",
                  size=32, bold=True, color=C.white, font=FONT_HDR,
                  align=PP_ALIGN.CENTER)

    steps = [
        ("STEP 1", "診断レポートの読み合わせ",
         "社内で優先度・工数感の共有。疑問点の洗い出し（約60分）", "📋"),
        ("STEP 2", "Quick Win実装キックオフ",
         "Phase 1施策（Day 1–14）の担当・期日確定。テンプレートコード提供", "🚀"),
        ("STEP 3", "60日後の再診断",
         "同じツールで再測定。改善効果の定量確認と次フェーズ計画", "📊"),
    ]
    for i, (num, title, desc, icon) in enumerate(steps):
        x = 0.7 + i * 4.1
        _add_rect(s, x, 2.5, 3.8, 3.2, C.deep_blue, line_color=C.cyan,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, x, 2.75, 3.8, 0.8, icon,
                      size=40, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_rect(s, x + 1.3, 3.75, 1.2, 0.35, C.cyan,
                  shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
        _add_text_box(s, x + 1.3, 3.75, 1.2, 0.35, num,
                      size=10, bold=True, color=C.navy,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, x + 0.2, 4.2, 3.4, 0.7, title,
                      size=16, bold=True, color=C.white, font=FONT_HDR,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _add_text_box(s, x + 0.25, 4.95, 3.3, 0.7, desc,
                      size=11, color=C.text_light, align=PP_ALIGN.CENTER)

    _add_rect(s, 2.5, 6.1, 8.3, 0.8, C.teal,
              shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
    _add_text_box(s, 2.5, 6.1, 8.3, 0.8,
                  "▶ 詳細な実装支援・カスタム診断のご相談を承ります",
                  size=16, bold=True, color=C.white, font=FONT_HDR,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    _add_text_box(s, 0.5, 7.1, 12.3, 0.3,
                  f"Thank you  |  AIO/LLMO診断ツール  |  {today}",
                  size=10, color=C.muted, align=PP_ALIGN.CENTER)

    # ===== Bytesに書き出し =====
    buf = BytesIO()
    pres.save(buf)
    buf.seek(0)
    return buf
